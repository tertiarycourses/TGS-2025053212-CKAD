"""
Inspect PPTX image slides and match them against source PDFs.
Shows consecutive same-title runs and flags visual similarity.

Usage:
    python inspect_slides.py <deck.pptx> [--range 40-60]
"""
import sys, os, hashlib, argparse
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SOURCE_PDFS = {
    "Day1":  BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf",
    "Day2":  BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf",
    "Day3":  BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf",
    "Day4M": BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf",
}
RENDER_DPI = 2.0


def build_pdf_index():
    index = {}
    for pdf_name, path in SOURCE_PDFS.items():
        if not os.path.exists(path):
            continue
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(RENDER_DPI, RENDER_DPI)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            blob = pix.tobytes("png")
            h = hashlib.md5(blob).hexdigest()
            lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
            title = lines[0][:70] if lines else "(blank)"
            index[h] = (pdf_name, i + 1, title)
        doc.close()
    return index


def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def get_img_blob(slide):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh.image.blob
    return None

def slide_text(slide):
    return " | ".join(
        sh.text_frame.text.strip()
        for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    )


def inspect(pptx_path, slide_range=None):
    print(f"Building PDF index...", flush=True)
    idx = build_pdf_index()
    print(f"Indexed {len(idx)} PDF pages\n")

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    total = len(slides)

    start, end = (1, total)
    if slide_range:
        parts = slide_range.split('-')
        start = int(parts[0])
        end = int(parts[1]) if len(parts) > 1 else start

    print(f"{'Slide':<6} {'Type':<10} {'Source':<22} {'Title'}")
    print("-" * 90)

    prev_title = None
    run_count = 0

    for i, slide in enumerate(slides):
        num = i + 1
        if num < start or num > end:
            continue

        if is_fullslide_image(slide, prs):
            blob = get_img_blob(slide)
            h = hashlib.md5(blob).hexdigest() if blob else "no-blob"
            if h in idx:
                pdf_name, page_num, title = idx[h]
                source = f"{pdf_name} p{page_num}"
            else:
                source = "(unknown)"
                title = "(not in source PDFs)"

            # Track same-title runs
            if title == prev_title:
                run_count += 1
                flag = f"  ← same title x{run_count+1}" if run_count > 0 else ""
            else:
                run_count = 0
                flag = ""

            print(f"{num:<6} {'IMAGE':<10} {source:<22} {title}{flag}")
            prev_title = title
        else:
            txt = slide_text(slide)[:60]
            print(f"{num:<6} {'TEXT':<10} {'—':<22} {txt}")
            prev_title = None
            run_count = 0

    print("\n")
    # Report all same-title runs across full deck
    print("Same-title image slide runs in full deck:")
    print("-" * 60)
    title_positions = {}
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        h = hashlib.md5(blob).hexdigest() if blob else ""
        title = idx[h][2] if h in idx else "(unknown)"
        if title not in title_positions:
            title_positions[title] = []
        title_positions[title].append(i + 1)

    found_runs = False
    for title, positions in sorted(title_positions.items(), key=lambda x: x[1][0]):
        if len(positions) >= 2:
            found_runs = True
            print(f"  [{len(positions)}x] slides {positions}: \"{title}\"")

    if not found_runs:
        print("  No same-title runs found — deck is clean.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input", help="PPTX path")
    parser.add_argument("--range", help="Slide range e.g. 45-60")
    args = parser.parse_args()
    inspect(args.input, args.range)
