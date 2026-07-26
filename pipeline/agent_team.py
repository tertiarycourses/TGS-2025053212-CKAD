"""
agent_team.py — CKAD Slide Pipeline Agent Team
================================================
Four agents with clear responsibilities:

  ImportAgent  — copies admin slides from v2.0, runs fix_zip_dups
  AuditAgent   — runs audit_ckad.py, returns structured violations
  FixAgent     — reads violations, applies targeted fixes
  Orchestrator — wires them together; loops until PASS or max iterations

Usage:
    python pipeline/agent_team.py [--deck <path/to/deck.pptx>] [--max-iter N]

Defaults:
    --deck      courseware/CKAD-Certified-Kubernetes-Application-Developer-v4.5-admin2.pptx
    --max-iter  5
"""

import sys, os, json, copy, re, argparse, subprocess, time
from io import BytesIO
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
PIPELINE  = os.path.join(BASE_DIR, "pipeline")
V20_PATH  = os.path.join(BASE_DIR, "courseware",
            "CKAD-Certified-Kubernetes-Application-Developer-v2.0.pptx")

ADMIN_MAP = {0: 1, 1: 2, 2: 3, 3: 7, 4: 9, 5: 10}  # dst slot → v2.0 src index

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# ─────────────────────────────────────────────────────────────────────────────
# Helpers shared between agents
# ─────────────────────────────────────────────────────────────────────────────

def log(agent, msg):
    print(f"  [{agent}] {msg}", flush=True)


def run_py(script, *args):
    """Run a pipeline Python script as a subprocess; return (returncode, stdout)."""
    result = subprocess.run(
        [sys.executable, os.path.join(PIPELINE, script), *args],
        capture_output=True, text=True, encoding='utf-8'
    )
    return result.returncode, result.stdout + result.stderr


# ─────────────────────────────────────────────────────────────────────────────
# Agent 1 — ImportAgent
# ─────────────────────────────────────────────────────────────────────────────

class ImportAgent:
    """
    Responsibility: ensure the first 6 slides are the correct admin slides
    from v2.0, with images properly transferred.
    """
    NAME = "ImportAgent"

    def run(self, deck_path: str) -> dict:
        log(self.NAME, f"Checking admin slides in {os.path.basename(deck_path)}")
        prs = Presentation(deck_path)
        slides = list(prs.slides)

        if len(slides) < 6:
            return {"status": "error", "message": "Deck has fewer than 6 slides"}

        # Check if first 6 look like admin slides already
        need_import = False
        admin_keywords = ["COURSE ADMINISTRATION", "TRAQOM", "YOUR TRAINER",
                          "SCHEDULE", "FINAL ASSESSMENT", "ASSESSMENT RULES"]
        for i, kw in enumerate(admin_keywords):
            txt = " ".join(
                sh.text_frame.text.upper()
                for sh in slides[i].shapes if sh.has_text_frame
            )
            if kw not in txt:
                need_import = True
                log(self.NAME, f"  Slide {i+1} missing expected content: {kw}")

        if not need_import:
            log(self.NAME, "Admin slides already present — skip import.")
            return {"status": "skipped", "message": "Admin slides OK"}

        log(self.NAME, "Importing admin slides from v2.0 ...")
        src = Presentation(V20_PATH)

        for dst_idx, src_idx in sorted(ADMIN_MAP.items()):
            src_slide = src.slides[src_idx]
            dst_slide = prs.slides[dst_idx]
            self._replace_content(src_slide, dst_slide)
            label = " | ".join(
                sh.text_frame.text.strip()[:40]
                for sh in src_slide.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()
            )[:60]
            log(self.NAME, f"  Slot {dst_idx+1} <- v2.0[{src_idx+1}]: {label}")

        prs.save(deck_path)
        log(self.NAME, "Saved. Running fix_zip_dups ...")
        rc, out = run_py("fix_zip_dups.py", deck_path)
        log(self.NAME, out.strip().splitlines()[-1] if out.strip() else "done")
        return {"status": "imported", "message": "Admin slides replaced from v2.0"}

    def _replace_content(self, src_slide, dst_slide):
        rId_map = {}
        for rId, rel in src_slide.part.rels.items():
            if rel.is_external or "/image" not in rel.reltype:
                continue
            try:
                blob = rel.target_part.blob
                _, new_rId = dst_slide.part.get_or_add_image_part(BytesIO(blob))
                rId_map[rId] = new_rId
            except Exception as e:
                log(self.NAME, f"  [WARN] image {rId}: {e}")
        dst_sp = dst_slide.shapes._spTree
        src_sp = src_slide.shapes._spTree
        for ch in list(dst_sp):
            dst_sp.remove(ch)
        for ch in src_sp:
            new_ch = copy.deepcopy(ch)
            for el in new_ch.iter():
                for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
                    v = el.get(attr)
                    if v and v in rId_map:
                        el.set(attr, rId_map[v])
            dst_sp.append(new_ch)


# ─────────────────────────────────────────────────────────────────────────────
# Agent 2 — AuditAgent
# ─────────────────────────────────────────────────────────────────────────────

class AuditAgent:
    """
    Responsibility: run audit_ckad.py and return structured violations.
    """
    NAME = "AuditAgent"

    def run(self, deck_path: str) -> dict:
        log(self.NAME, f"Auditing {os.path.basename(deck_path)} ...")
        rc, out = run_py("audit_ckad.py", deck_path)

        # Parse JSON from output
        try:
            # audit_ckad prints JSON to stdout then a summary to stderr
            json_str = out[out.index('{'):out.rindex('}')+1]
            report = json.loads(json_str)
        except Exception:
            return {"status": "error", "raw": out}

        passed = report.get("pass", False)
        violations = report.get("violations", [])

        if passed:
            log(self.NAME, f"PASS — {report['total_slides']} slides, "
                           f"{len(report['labs_found'])} labs")
        else:
            log(self.NAME, f"FAIL — {len(violations)} violation(s):")
            for v in violations:
                log(self.NAME, f"   slide {v['slide']} | {v['rule']} | {v.get('detail','')[:60]}")

        return {
            "status": "pass" if passed else "fail",
            "violations": violations,
            "report": report
        }


# ─────────────────────────────────────────────────────────────────────────────
# Agent 3 — FixAgent
# ─────────────────────────────────────────────────────────────────────────────

class FixAgent:
    """
    Responsibility: given a list of audit violations, apply targeted fixes.
    Returns a list of actions taken.
    """
    NAME = "FixAgent"

    FIXABLE = {
        "admin_slide",    # re-import triggers ImportAgent
        "domain_order",   # re-run reorder
        "missing_labs",   # flag only — cannot auto-fix
        "missing_killercoda_url",  # flag only
    }

    def run(self, deck_path: str, violations: list) -> dict:
        if not violations:
            return {"status": "nothing_to_fix", "actions": []}

        rules = {v["rule"] for v in violations}
        actions = []

        log(self.NAME, f"Fixing: {', '.join(rules)}")

        # Fix domain ordering by re-running reorder_ckad
        if "domain_order" in rules:
            log(self.NAME, "Re-running reorder_ckad.py ...")
            rc, out = run_py("reorder_ckad.py", deck_path, deck_path)
            log(self.NAME, out.strip().splitlines()[-1] if out.strip() else "done")
            actions.append("reorder_ckad")

        # Admin slides out of place — signal ImportAgent to re-run
        if "admin_slide" in rules:
            log(self.NAME, "Admin slide violations — flagging for ImportAgent re-run")
            actions.append("reimport_admin")

        # Unfixable — report only
        for rule in rules - {"admin_slide", "domain_order"}:
            log(self.NAME, f"  [MANUAL] Cannot auto-fix: {rule}")
            actions.append(f"manual:{rule}")

        # Re-run fix_zip_dups after any changes
        if {"domain_order", "admin_slide"} & rules:
            log(self.NAME, "Running fix_zip_dups ...")
            rc, out = run_py("fix_zip_dups.py", deck_path)
            log(self.NAME, out.strip().splitlines()[-1] if out.strip() else "done")

        return {"status": "fixed", "actions": actions}


# ─────────────────────────────────────────────────────────────────────────────
# Orchestrator
# ─────────────────────────────────────────────────────────────────────────────

class Orchestrator:
    """
    Wires ImportAgent → AuditAgent → (FixAgent → AuditAgent) loop.
    """
    def __init__(self, deck_path: str, max_iter: int = 5):
        self.deck_path = deck_path
        self.max_iter  = max_iter
        self.import_agent = ImportAgent()
        self.audit_agent  = AuditAgent()
        self.fix_agent    = FixAgent()

    def run(self) -> dict:
        start = time.time()
        print(f"\n{'='*60}")
        print(f"  CKAD AGENT TEAM — {os.path.basename(self.deck_path)}")
        print(f"{'='*60}\n")

        # ── Step 1: Import ────────────────────────────────────────
        print("── Step 1: ImportAgent ──────────────────────────────────")
        import_result = self.import_agent.run(self.deck_path)
        print()

        # ── Step 2: Audit → Fix loop ──────────────────────────────
        history = []
        for iteration in range(1, self.max_iter + 1):
            print(f"── Step 2 (iteration {iteration}/{self.max_iter}): AuditAgent ──────────────")
            audit_result = self.audit_agent.run(self.deck_path)
            history.append(audit_result)
            print()

            if audit_result["status"] == "pass":
                break

            if iteration == self.max_iter:
                print(f"── Max iterations reached — stopping.\n")
                break

            print(f"── Step 3 (iteration {iteration}): FixAgent ─────────────────────")
            fix_result = self.fix_agent.run(self.deck_path, audit_result["violations"])
            print()

            # If FixAgent flagged admin re-import, run ImportAgent again
            if "reimport_admin" in fix_result.get("actions", []):
                print("── ImportAgent re-run (triggered by FixAgent) ───────────")
                self.import_agent.run(self.deck_path)
                print()

        # ── Final summary ─────────────────────────────────────────
        elapsed = time.time() - start
        final   = history[-1] if history else {}
        passed  = final.get("status") == "pass"

        print(f"{'='*60}")
        print(f"  RESULT : {'✅ PASS' if passed else '❌ FAIL'}")
        print(f"  Slides : {final.get('report', {}).get('total_slides', '?')}")
        print(f"  Labs   : {len(final.get('report', {}).get('labs_found', []))}/30")
        print(f"  Iters  : {len(history)}")
        print(f"  Time   : {elapsed:.1f}s")
        if not passed:
            remaining = final.get("violations", [])
            print(f"  Remaining violations ({len(remaining)}):")
            for v in remaining:
                print(f"    slide {v['slide']} | {v['rule']} | {v.get('detail','')[:60]}")
        print(f"{'='*60}\n")

        return {
            "pass": passed,
            "iterations": len(history),
            "import": import_result,
            "final_audit": final,
        }


# ─────────────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="CKAD Agent Team Pipeline")
    ap.add_argument(
        "--deck", default=os.path.join(
            BASE_DIR, "courseware",
            "CKAD-Certified-Kubernetes-Application-Developer-v4.5-admin2.pptx"
        ),
        help="Path to the PPTX deck to process"
    )
    ap.add_argument("--max-iter", type=int, default=5,
                    help="Maximum audit→fix iterations")
    args = ap.parse_args()

    orch = Orchestrator(deck_path=args.deck, max_iter=args.max_iter)
    result = orch.run()
    sys.exit(0 if result["pass"] else 1)
