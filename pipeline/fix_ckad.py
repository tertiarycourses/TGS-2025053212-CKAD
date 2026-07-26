"""
CKAD Slide Fixer
================
Reads an audit JSON report (from audit_ckad.py) and fixes the violations
it can handle automatically:

  - admin_slide          → remove the slide
  - missing_killercoda_url → add URL text box to lab header slide

Violations it cannot fix (require human/importer):
  - missing_labs         → flagged, not fixed
  - domain_order         → flagged, not fixed

Usage:
    python audit_ckad.py deck.pptx > report.json
    python fix_ckad.py deck.pptx report.json fixed_deck.pptx
"""
import sys, re, json, io
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

C_BLUE = RGBColor(0x02, 0x9A, 0xE8)

LAB_URLS = {
    1: "https://killercoda.com/playgrounds/scenario/ubuntu",
    2: "https://killercoda.com/playgrounds/scenario/ubuntu",
}
for n in range(3, 31):
    LAB_URLS[n] = "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/two-node"


def slide_text(slide):
    return " | ".join(
        sh.text_frame.text.strip()
        for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    )

def lab_num(text_upper):
    m = re.search(r'\bLAB\s+(\d+)\b', text_upper)
    return int(m.group(1)) if m else None

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def add_url(slide, url):
    txBox = slide.shapes.add_textbox(Inches(0.85), Inches(6.45), Inches(11.70), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = url
    run.font.name   = "Calibri"
    run.font.size   = Pt(11)
    run.font.italic = True
    run.font.color.rgb = C_BLUE
    run.font.bold   = False


def fix(pptx_path, report_path, out_path):
    with open(report_path, encoding='utf-8') as f:
        report = json.load(f)

    if report["pass"]:
        print("Nothing to fix — audit already passes.")
        import shutil; shutil.copy(pptx_path, out_path)
        return

    violations = report["violations"]
    admin_slides  = {v["slide"] for v in violations if v["rule"] == "admin_slide"}
    url_slides    = {v["slide"]: v["detail"] for v in violations if v["rule"] == "missing_killercoda_url"}
    unfixable     = [v for v in violations if v["rule"] in ("missing_labs", "domain_order")]

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    print(f"Loaded: {len(slides)} slides")

    # Fix missing URLs first (before reindexing from removals)
    for slide_num, detail in url_slides.items():
        i = slide_num - 1
        if i >= len(slides): continue
        slide = slides[i]
        up = slide_text(slide).upper()
        n = lab_num(up)
        if n and n in LAB_URLS:
            add_url(slide, LAB_URLS[n])
            print(f"  + URL added → slide {slide_num} Lab {n}")

    # Remove admin slides (by index, remove from highest to avoid reindex)
    remove_idx = {s - 1 for s in admin_slides if 1 <= s <= len(slides)}
    print(f"  - Removing {len(remove_idx)} admin slides: {sorted(s+1 for s in remove_idx)}")
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    kept = [el for j, el in enumerate(all_ids) if j not in remove_idx]
    for el in list(sp): sp.remove(el)
    for el in kept: sp.append(el)

    if unfixable:
        print("\nCannot auto-fix:")
        for v in unfixable:
            print(f"  [{v['rule']}] {v['detail']}")

    print(f"\nSaving → {out_path}")
    prs.save(out_path)
    print(f"Done! {len(kept)} slides")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python fix_ckad.py <deck.pptx> <report.json> <output.pptx>")
        sys.exit(2)
    fix(sys.argv[1], sys.argv[2], sys.argv[3])
