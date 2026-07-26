"""
CKAD Slide Deduplication Agent  (v2)
=====================================
Two-pass deduplication:

  Pass 1 — Exact pixel duplicates (MD5 hash of image bytes).
            Keeps first occurrence, removes all later exact copies.

  Pass 2 — Source-PDF fingerprint: builds a hash index of every page
            from the 5 source PDFs (rendered at build DPI) and removes
            any PPTX image slide whose bytes match a PDF page that was
            already added earlier in the deck.
            This catches same-content slides that may have been re-added
            by overlapping section ranges.

  Pass 3 — Same-title consecutive dedupe: if N consecutive image slides
            all have the same title text (extracted via OCR/fitz), keep
            only the first occurrence. (Skipped when fitz OCR unavailable.)

Usage:
    python dedupe_ckad.py <input.pptx> [output.pptx] [--loop]

  --loop: keeps re-running passes until no changes are made (convergence).

Exit 0 always. Prints JSON summary to stdout.
"""
import sys, os, hashlib, json, re, argparse, io
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SOURCE_PDFS = {
    "Day1":  BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf",
    "Day2":  BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf",
    "Day3":  BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf",
    "Day4M": BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf",
}
RENDER_DPI = 2.0


# ── helpers ─────────────────────────────────────────────────────────────────

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def get_img_blob(slide):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh.image.blob
    return None

def build_pdf_index():
    """MD5 of every rendered PDF page → (pdf_name, page_num, title)."""
    index = {}
    for pdf_name, path in SOURCE_PDFS.items():
        if not os.path.exists(path):
            continue
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(RENDER_DPI, RENDER_DPI)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            blob = pix.tobytes("png")
            h = hashlib.md5(blob).hexdigest()
            lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
            title = lines[0][:60] if lines else "(blank)"
            index[h] = (pdf_name, i + 1, title)
        doc.close()
    return index

def extract_title_from_blob(blob):
    """Try to extract the first non-empty text line from an image slide."""
    try:
        doc = fitz.open(stream=blob, filetype="png")
        page = doc[0]
        lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
        return lines[0][:80] if lines else ""
    except Exception:
        return ""

def remove_slides(prs, remove_idx):
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    kept = [el for j, el in enumerate(all_ids) if j not in remove_idx]
    for el in list(sp):
        sp.remove(el)
    for el in kept:
        sp.append(el)
    return len(kept)


# ── passes ───────────────────────────────────────────────────────────────────

def pass1_exact(prs):
    """Remove exact pixel duplicates (MD5)."""
    slides = list(prs.slides)
    seen, remove = {}, set()
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in seen:
            remove.add(i)
        else:
            seen[h] = i + 1
    return remove


def pass2_pdf_source(prs, pdf_index):
    """Remove image slides whose source PDF page was already injected earlier."""
    slides = list(prs.slides)
    seen_pages, remove = set(), set()
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in pdf_index:
            page_key = (pdf_index[h][0], pdf_index[h][1])  # (pdf_name, page_num)
            if page_key in seen_pages:
                remove.add(i)
            else:
                seen_pages.add(page_key)
    return remove


def pass3_title_global(prs, pdf_index, max_per_title=1):
    """
    Global title dedup: for each unique title (from PDF index), keep only
    the first `max_per_title` occurrences in the deck.
    Skips slides whose source is unknown (not in pdf_index).
    """
    slides = list(prs.slides)
    title_seen = {}   # title → list of (slide_idx,)
    remove = set()

    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h not in pdf_index:
            continue  # unknown origin — don't touch
        title = pdf_index[h][2]
        if not title or title == "(blank)":
            continue

        if title not in title_seen:
            title_seen[title] = 0
        title_seen[title] += 1
        if title_seen[title] > max_per_title:
            remove.add(i)

    return remove


# ── main ─────────────────────────────────────────────────────────────────────

def dedupe_once(prs, pdf_index, max_per_title=None):
    """Run all passes on prs (in-place). Returns total removed."""
    total_removed = 0
    removed_slides = []

    # Pass 1 — exact pixel duplicates
    r1 = pass1_exact(prs)
    if r1:
        removed_slides += [i + 1 for i in sorted(r1)]
        total_removed += len(r1)
        remove_slides(prs, r1)

    # Pass 2 — same source PDF page appearing twice
    r2 = pass2_pdf_source(prs, pdf_index)
    if r2:
        removed_slides += [i + 1 for i in sorted(r2)]
        total_removed += len(r2)
        remove_slides(prs, r2)

    # Pass 3 — title-based global dedup (optional)
    if max_per_title is not None:
        r3 = pass3_title_global(prs, pdf_index, max_per_title)
        if r3:
            removed_slides += [i + 1 for i in sorted(r3)]
            total_removed += len(r3)
            remove_slides(prs, r3)

    return total_removed, removed_slides


def dedupe(input_path, output_path=None, loop=False, max_per_title=None):
    if output_path is None:
        output_path = input_path

    print("  Building PDF source index...", file=sys.stderr)
    pdf_index = build_pdf_index()
    print(f"  Indexed {len(pdf_index)} PDF pages", file=sys.stderr)

    prs = Presentation(input_path)
    total_in = len(list(prs.slides))
    all_removed = []
    iteration = 0

    while True:
        iteration += 1
        n, removed = dedupe_once(prs, pdf_index, max_per_title)
        all_removed.extend(removed)
        print(f"  Iteration {iteration}: removed {n} slides", file=sys.stderr)
        if not loop or n == 0:
            break

    if all_removed or output_path != input_path:
        prs.save(output_path)

    total_out = len(list(prs.slides))
    result = {
        "input":              input_path,
        "output":             output_path,
        "total_in":           total_in,
        "total_out":          total_out,
        "iterations":         iteration,
        "duplicates_removed": len(all_removed),
        "removed_slides":     sorted(all_removed),
        "clean":              len(all_removed) == 0,
    }
    return result


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input",   help="Input PPTX")
    parser.add_argument("output",  nargs="?", help="Output PPTX (default: overwrite input)")
    parser.add_argument("--loop",  action="store_true", help="Keep iterating until no changes")
    parser.add_argument("--max-per-title", type=int, default=None,
                        help="Max slides with same title to keep (e.g. 1 = keep only first)")
    args = parser.parse_args()

    out = args.output or args.input
    result = dedupe(args.input, out, loop=args.loop, max_per_title=args.max_per_title)
    print(json.dumps(result, indent=2))

    status = "CLEAN" if result["clean"] else f"REMOVED {result['duplicates_removed']}"
    print(f"\nDEDUPE {status} — {result['total_out']} slides remaining", file=sys.stderr)
    sys.exit(0)
