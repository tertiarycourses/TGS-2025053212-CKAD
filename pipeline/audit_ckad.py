"""
CKAD Slide Auditor
==================
Scans a PPTX and returns a JSON report of every quality violation.
Exit code 0 = PASS (no violations). Exit code 1 = FAIL.

Usage:
    python audit_ckad.py <path/to/deck.pptx>

Output (stdout, JSON):
    {
      "total_slides": 460,
      "violations": [
        {"slide": 12, "rule": "admin_slide", "detail": "tea break"},
        ...
      ],
      "pass": true
    }
"""
import sys, re, json
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Rules ──────────────────────────────────────────────────────────────────
ADMIN_PATTERNS = [
    "course slides", "traqom", "ssg digital attendance", "your trainer",
    "final assessment | assessment", "tea break", "lunch break", "wrap-up",
    "day 1 objectives", "day 2 objectives", "day 3 objectives", "day 4 objectives",
    "topic 1 |", "topic 2 |", "topic 3 |", "topic 4 |", "topic 5 |",
    "topic 6 |", "topic 7 |", "topic 8 |", "topic 9 |", "topic 10 |",
    "topic 11 |", "topic 12 |",
    "assessment rules", "warm-up | mock", "schedule | lesson plan",
    "house rules | ground rules", "course administration", "see you tomorrow",
    "thank you everyone", "basic ground rules", "topics we",
    "day 4  (½ day) | services",
]

# Expected labs in order
EXPECTED_LABS = list(range(1, 31))  # 1–30

LAB_URLS = {
    1: "killercoda.com/playgrounds/scenario/ubuntu",
    2: "killercoda.com/playgrounds/scenario/ubuntu",
}
for n in range(3, 31):
    LAB_URLS[n] = "killercoda.com/playgrounds/course/kubernetes-playgrounds/two-node"

# Expected domain headers (in order)
DOMAIN_MARKERS = [
    "DAY 1",     # Domain 1
    "DAY 2",     # Domain 2
    "DOMAIN 3",  # Domain 3
    "DOMAIN 4",  # Domain 4
    "DAY 4",     # Domain 5
]


# ── Helpers ─────────────────────────────────────────────────────────────────
def slide_text(slide):
    return " | ".join(
        sh.text_frame.text.strip()
        for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    )

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def lab_num(text_upper):
    m = re.search(r'\bLAB\s+(\d+)\b', text_upper)
    return int(m.group(1)) if m else None


# ── Audit ───────────────────────────────────────────────────────────────────
def audit(pptx_path):
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    violations = []

    seen_labs = []
    seen_domains = []

    for i, slide in enumerate(slides):
        snum = i + 1
        raw = slide_text(slide)
        low = raw.lower()
        up  = raw.upper()

        fullslide = is_fullslide_image(slide, prs)

        # Rule 1: no admin slides (even if they have small images)
        if not fullslide:
            wc = len(re.findall(r'[a-zA-Z]{3,}', low))
            if wc >= 2:
                for p in ADMIN_PATTERNS:
                    if p in low:
                        violations.append({
                            "slide": snum, "rule": "admin_slide",
                            "detail": p,
                            "text": raw[:80]
                        })
                        break

        # Rule 2: every lab header must have KillerCoda URL
        if not fullslide and re.search(r'\bLAB\s+\d+\b', up):
            if "DAY" in up or "DOMAIN" in up:
                n = lab_num(up)
                if n:
                    seen_labs.append((snum, n))
                    expected_url = LAB_URLS.get(n, "")
                    if expected_url and expected_url not in low:
                        violations.append({
                            "slide": snum, "rule": "missing_killercoda_url",
                            "detail": f"Lab {n} — expected {expected_url}",
                            "text": raw[:80]
                        })

        # Rule 3: track domain markers for ordering check
        for dm in DOMAIN_MARKERS:
            if dm in up and dm not in seen_domains:
                seen_domains.append(dm)

    # Rule 4: all 30 labs must be present
    found_lab_nums = sorted(set(n for _, n in seen_labs))
    missing_labs = [n for n in EXPECTED_LABS if n not in found_lab_nums]
    if missing_labs:
        violations.append({
            "slide": 0, "rule": "missing_labs",
            "detail": f"Labs not found: {missing_labs}"
        })

    # Rule 5: domain markers must appear in correct order
    expected_order = [dm for dm in DOMAIN_MARKERS if dm in seen_domains]
    if expected_order != seen_domains:
        violations.append({
            "slide": 0, "rule": "domain_order",
            "detail": f"Expected {DOMAIN_MARKERS}, found {seen_domains}"
        })

    result = {
        "file": pptx_path,
        "total_slides": len(slides),
        "labs_found": found_lab_nums,
        "domains_found": seen_domains,
        "violations": violations,
        "pass": len(violations) == 0,
    }
    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python audit_ckad.py <deck.pptx>")
        sys.exit(2)

    report = audit(sys.argv[1])
    print(json.dumps(report, indent=2))

    if report["pass"]:
        print(f"\nAUDIT PASS — {report['total_slides']} slides, all {len(report['labs_found'])} labs present", file=sys.stderr)
        sys.exit(0)
    else:
        print(f"\nAUDIT FAIL — {len(report['violations'])} violations", file=sys.stderr)
        sys.exit(1)
