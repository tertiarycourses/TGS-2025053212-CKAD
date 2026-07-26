"""
insert_admin_slides.py  (v2 — proper image copy)
=================================================
Copies the admin/intro slides from v2.0 and inserts them at the beginning
of v4.5.  Properly transfers image parts (icon PNGs/SVGs) so shapes with
pictures are not broken.

Slides copied from v2.0 (0-indexed):
  idx 1  — Welcome & Housekeeping
  idx 2  — Digital Attendance (Mandatory)
  idx 3  — About the Trainer (Mohan Pothula)
  idx 7  — Lesson Plan / Schedule
  idx 9  — Final Assessment
  idx 10 — Assessment Rules

The v2.0 slides already use the correct design (white bg + small accent +
dark text) — no text-colour or left-bar changes are applied.

Usage:
    python pipeline/insert_admin_slides.py
"""
import sys, copy
from io import BytesIO
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SRC_V20  = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v2.0.pptx"
IN_V45   = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v4.5-final.pptx"
OUT_V45  = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v4.5-final.pptx"

# v2.0 slide indices (0-based) to copy, in order
SLIDES_TO_COPY = [1, 2, 3, 7, 9, 10]

# r: namespace for rId references in shape XML
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def copy_slide_with_images(src_prs, src_idx, dst_prs):
    """
    Copy slide src_idx from src_prs to the END of dst_prs.
    Properly transfers PICTURE image parts so icons are not broken.
    Text colours and bold settings are preserved as-is from v2.0.
    Returns the new slide.
    """
    src_slide = src_prs.slides[src_idx]

    # Step 1: Add blank slide to destination
    blank = dst_prs.slide_layouts[6]
    new_slide = dst_prs.slides.add_slide(blank)

    # Step 2: Transfer image parts from source slide to new slide
    # Build old-rId → new-rId map
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.is_external:
            continue
        if "/image" in rel.reltype:
            blob = rel.target_part.blob
            try:
                # get_or_add_image_part accepts a file-like object (BytesIO)
                _img_part, new_rId = new_slide.part.get_or_add_image_part(BytesIO(blob))
                rId_map[rId] = new_rId
                print(f"      image rId {rId} -> {new_rId}")
            except Exception as e:
                print(f"      [WARN] could not copy image rId {rId}: {e}")

    # Step 3: Deep-copy spTree children, rewriting rId references
    dst_sp = new_slide.shapes._spTree
    src_sp = src_slide.shapes._spTree

    # Clear the blank slide's default placeholder shapes
    for ch in list(dst_sp):
        dst_sp.remove(ch)

    for ch in src_sp:
        new_ch = copy.deepcopy(ch)
        # Rewrite every r:embed and r:link attribute in the copied XML
        for el in new_ch.iter():
            for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
                old_val = el.get(attr)
                if old_val and old_val in rId_map:
                    el.set(attr, rId_map[old_val])
        dst_sp.append(new_ch)

    return new_slide


def move_slides_to_front(prs, count):
    """Move the last `count` slides (just appended) to positions 0..count-1."""
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    front   = all_ids[-count:]
    rest    = all_ids[:-count]
    for el in list(sp):
        sp.remove(el)
    for el in (front + rest):
        sp.append(el)


def slide_label(prs, idx):
    sl = prs.slides[idx]
    parts = []
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()[:50]
            if t:
                parts.append(t)
    return " | ".join(parts[:2])


def main():
    print("Opening source presentations...")
    src = Presentation(SRC_V20)
    dst = Presentation(IN_V45)
    print(f"  v2.0: {len(list(src.slides))} slides")
    print(f"  v4.5: {len(list(dst.slides))} slides (before insert)")

    # Remove any previously-inserted admin slides (first 6) so we don't double-insert.
    # Detection: v2.0 admin slides all have a white full-slide background AND no full-slide image.
    # Simple heuristic: if slide_count > expected, skip removal (let user decide).
    # For safety: just re-insert — caller should start from a clean v4.5 without admin slides.

    print(f"\nCopying {len(SLIDES_TO_COPY)} admin slides from v2.0 ...")
    for idx in SLIDES_TO_COPY:
        label = slide_label(src, idx)
        print(f"  v2.0[{idx+1}]: {label}")
        copy_slide_with_images(src, idx, dst)

    move_slides_to_front(dst, len(SLIDES_TO_COPY))
    print(f"\nMoved {len(SLIDES_TO_COPY)} slides to front of deck.")

    dst.save(OUT_V45)
    total = len(list(dst.slides))
    print(f"Saved -> {OUT_V45}  ({total} slides total)")


if __name__ == "__main__":
    main()
