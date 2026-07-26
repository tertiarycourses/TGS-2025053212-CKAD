"""
CKAD Slide Comprehensive Validator
====================================
Runs all quality checks in one pass:

  Check 1 — Lab completeness:  all 30 labs must be present.
  Check 2 — KillerCoda URLs:   every lab header slide must have the correct URL.
  Check 3 — No admin slides:   no housekeeping / break slides.
  Check 4 — No exact duplicates (MD5): image slides must all be unique.
  Check 5 — No image-title duplicates: each PDF source title appears at most once.
  Check 6 — No text-vs-image topic overlap: text slides covered by an image slide.
  Check 7 — Content before labs: for each section, theory images must precede
             the first lab header slide.

Usage:
    python validate_ckad.py <deck.pptx> [--fix-output output.pptx]

  --fix-output  If supplied, applies auto-fixes for checks 4-6 (duplicate removal)
                and saves to the given path.  Checks 1-3 and 7 are report-only.

Exit codes:
  0 = all checks pass
  1 = violations found
"""
import sys, os, hashlib, json, re, argparse
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SOURCE_PDFS = {
    "Day1":  BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf",
    "Day2":  BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf",
    "Day3":  BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf",
    "Day4M": BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf",
    "Day4A": BASE_DIR + r"\CKAD Day 4 Afternoon - Summary & Tips.pdf",
}
RENDER_DPI = 2.0

EXPECTED_LABS = list(range(1, 31))

LAB_URLS = {1: "killercoda.com/playgrounds/scenario/ubuntu",
            2: "killercoda.com/playgrounds/scenario/ubuntu"}
for n in range(3, 31):
    LAB_URLS[n] = "killercoda.com/playgrounds/course/kubernetes-playgrounds/two-node"

ADMIN_PATTERNS = [
    "course slides", "traqom", "ssg digital attendance", "your trainer",
    "tea break", "lunch break", "wrap-up", "day 1 objectives", "day 2 objectives",
    "day 3 objectives", "day 4 objectives", "assessment rules", "warm-up | mock",
    "house rules | ground rules", "course administration", "see you tomorrow",
    "thank you everyone", "basic ground rules", "schedule | lesson plan",
    "final assessment | assessment",
]

TEXT_TITLE_NOISE = [
    r"^(DOMAIN\s+\d+\s*[·:—\-]+\s*)", r"^(DAY\s+\d+\s*[·:—\-]+\s*)",
    r"^(FOUNDATIONS\s*[|·:—\-]+\s*)", r"^(THEORY\s*[|·:—\-]+\s*)",
    r"^(CONCEPTS?\s*[|·:—\-]+\s*)", r"^(OVERVIEW\s*[|·:—\-]+\s*)",
    r"^(SECTION\s+\d+\s*[|·:—\-]+\s*)",
]

# PDF source ranges per section (mirrors reorder_ckad.py SECTIONS list)
# Used by check7 to distinguish images belonging to THIS section vs later sections
SECTIONS_PDF = [
    {"name": "Container Images",        "pdf": "Day1",  "p_start": 0,   "p_end": 28},
    {"name": "Pods",                    "pdf": "Day1",  "p_start": 28,  "p_end": 82},
    {"name": "Jobs",                    "pdf": "Day1",  "p_start": 82,  "p_end": 97},
    {"name": "CronJobs",               "pdf": "Day1",  "p_start": 97,  "p_end": 106},
    {"name": "Multi-container Pods",   "pdf": "Day3",  "p_start": 98,  "p_end": 113},
    {"name": "Volumes",                "pdf": "Day3",  "p_start": 10,  "p_end": 33},
    {"name": "Labels and Annotations", "pdf": "Day1",  "p_start": 106, "p_end": 120},
    {"name": "Deployments",            "pdf": "Day1",  "p_start": 122, "p_end": 136},
    {"name": "Rolling Updates",        "pdf": "Day1",  "p_start": 136, "p_end": 155},
    {"name": "Blue/Green and Canary",  "pdf": "Day1",  "p_start": 157, "p_end": 164},
    {"name": "HPA (Autoscaler)",       "pdf": "Day1",  "p_start": 164, "p_end": 170},
    {"name": "Helm",                   "pdf": "Day3",  "p_start": 34,  "p_end": 48},
    {"name": "Kustomize",              "pdf": "Day3",  "p_start": 48,  "p_end": 62},
    {"name": "DaemonSets/StatefulSets","pdf": "Day4M", "p_start": 9,   "p_end": 20},
    {"name": "Probes",                 "pdf": "Day3",  "p_start": 80,  "p_end": 92},
    {"name": "Metrics Server",         "pdf": "Day3",  "p_start": 92,  "p_end": 98},
    {"name": "API Deprecations",       "pdf": "Day3",  "p_start": 72,  "p_end": 80},
    {"name": "ConfigMaps",             "pdf": "Day2",  "p_start": 130, "p_end": 141},
    {"name": "Secrets",                "pdf": "Day2",  "p_start": 141, "p_end": 152},
    {"name": "SecurityContext",        "pdf": "Day2",  "p_start": 153, "p_end": 161},
    {"name": "ServiceAccounts",        "pdf": "Day2",  "p_start": 100, "p_end": 107},
    {"name": "RBAC",                   "pdf": "Day2",  "p_start": 83,  "p_end": 100},
    {"name": "CRDs",                   "pdf": "Day3",  "p_start": 62,  "p_end": 72},
    {"name": "Quotas and Limits",      "pdf": "Day2",  "p_start": 107, "p_end": 129},
    {"name": "Taints and Node Affinity","pdf": "Day4M","p_start": 20,  "p_end": 34},
    {"name": "Services",               "pdf": "Day2",  "p_start": 22,  "p_end": 55},
    {"name": "Ingress and TLS",        "pdf": "Day2",  "p_start": 55,  "p_end": 67},
    {"name": "NetworkPolicy",          "pdf": "Day2",  "p_start": 67,  "p_end": 82},
]

# Sections where theory images should precede the first lab header
# each entry: (section_name, [keywords_identifying_first_lab_header])
SECTION_ORDER_CHECKS = [
    ("Container Images",       ["CONTAINER IMAGES", "LAB 1"]),
    ("Pods",                   ["PODS", "LAB 3"]),
    ("Jobs",                   ["JOBS", "LAB 4"]),
    ("CronJobs",               ["CRONJOBS", "LAB 5"]),
    ("Multi-container Pods",   ["MULTI-CONTAINER PODS", "LAB 6"]),
    ("Volumes",                ["VOLUMES", "LAB 8"]),
    ("Deployments",            ["DEPLOYMENTS", "LAB 9"]),
    ("Rolling Updates",        ["ROLLING UPDATES", "LAB 10"]),
    ("Blue/Green",             ["BLUE/GREEN", "LAB 11"]),
    ("Helm",                   ["HELM", "LAB 13"]),
    ("Kustomize",              ["KUSTOMIZE", "LAB 14"]),
    ("DaemonSets/StatefulSets",["DAEMON", "LAB 15"]),
    ("Probes",                 ["PROBES", "LAB 16"]),
    ("Metrics Server",         ["METRICS", "LAB 18"]),
    ("ConfigMaps",             ["CONFIGMAPS", "LAB 21"]),
    ("Secrets",                ["SECRETS", "LAB 22"]),
    ("SecurityContext",        ["SECURITYCONTEXT", "LAB 23"]),
    ("ServiceAccounts",        ["SERVICEACCOUNTS", "LAB 24"]),
    ("RBAC",                   ["RBAC", "LAB 25"]),
    ("Quotas and Limits",      ["QUOTAS", "LAB 26"]),
    ("Services",               ["SERVICES", "LAB 27"]),
    ("Ingress and TLS",        ["INGRESS", "LAB 29"]),
    ("NetworkPolicy",          ["NETWORKPOLICY", "LAB 30"]),
]


# ── helpers ─────────────────────────────────────────────────────────────────

def normalise(s):
    s = s.lower().strip()
    for pat in TEXT_TITLE_NOISE:
        s = re.sub(pat, '', s, flags=re.IGNORECASE)
    s = re.sub(r'[^a-z0-9 ]', ' ', s)
    return re.sub(r'\s+', ' ', s).strip()

def build_pdf_index():
    index = {}
    for pdf_name, path in SOURCE_PDFS.items():
        if not os.path.exists(path):
            print(f"  [WARN] PDF not found: {os.path.basename(path)}", file=sys.stderr)
            continue
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(RENDER_DPI, RENDER_DPI)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            blob = pix.tobytes("png")
            h = hashlib.md5(blob).hexdigest()
            lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
            title = lines[0][:80] if lines else ""
            index[h] = (pdf_name, i + 1, title)
        doc.close()
    return index

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def get_img_blob(slide):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh.image.blob
    return None

def slide_text_parts(slide):
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                parts.append(t)
    return parts

def slide_upper(slide):
    return " | ".join(slide_text_parts(slide)).upper()

def extract_text_slide_title(slide):
    parts = slide_text_parts(slide)
    if not parts:
        return None
    candidates = []
    for part in parts:
        for line in [l.strip() for l in part.split('\n') if l.strip()]:
            if len(line) < 8:
                continue
            if re.match(r'^(LAB\s+\d+|STEP\s+\d+|DAY\s+\d+\s*·)', line, re.I):
                continue
            if re.match(r'^(Tertiary|©|http)', line, re.I):
                continue
            candidates.append(line)
    for c in candidates:
        if len(normalise(c).split()) >= 3:
            return c
    return candidates[0] if candidates else None

def lab_num(text_upper):
    m = re.search(r'\bLAB\s+(\d+)\b', text_upper)
    return int(m.group(1)) if m else None

def remove_slides(prs, remove_idx):
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    kept = [el for j, el in enumerate(all_ids) if j not in remove_idx]
    for el in list(sp):
        sp.remove(el)
    for el in kept:
        sp.append(el)


# ── checks ───────────────────────────────────────────────────────────────────

def check1_lab_completeness(slides, prs):
    """All 30 labs must appear (header or commands slide)."""
    seen_labs = set()
    for slide in slides:
        if is_fullslide_image(slide, prs):
            continue
        up = slide_upper(slide)
        if re.search(r'\bLAB\s+\d+\b', up):
            is_lab = ("DAY" in up or "DOMAIN" in up or
                      re.search(r'\bCOMMANDS?\b', up) or
                      "TEST IT" in up or re.match(r'LAB\s+\d+', up.strip()))
            if is_lab:
                n = lab_num(up)
                if n:
                    seen_labs.add(n)
    missing = [n for n in EXPECTED_LABS if n not in seen_labs]
    flags = []
    if missing:
        flags.append({"check": "1_lab_completeness", "missing_labs": missing,
                      "found": sorted(seen_labs)})
    return flags


def check2_killercoda_urls(slides, prs):
    """Every lab header slide (with DAY/DOMAIN + LAB N) must have the correct URL."""
    flags = []
    for i, slide in enumerate(slides):
        if is_fullslide_image(slide, prs):
            continue
        up = slide_upper(slide)
        low = " ".join(slide_text_parts(slide)).lower()
        if re.search(r'\bLAB\s+\d+\b', up) and ("DAY" in up or "DOMAIN" in up):
            n = lab_num(up)
            if n:
                expected = LAB_URLS.get(n, "")
                if expected and expected not in low:
                    flags.append({"check": "2_killercoda_url", "slide": i+1,
                                  "lab": n, "expected": expected,
                                  "detail": f"Lab {n} — expected {expected}"})
    return flags


def check3_admin_slides(slides, prs):
    """No housekeeping / admin slides — skips intentional intro slides before first section divider."""
    flags = []
    content_started = False
    for i, slide in enumerate(slides):
        if is_fullslide_image(slide, prs):
            continue
        parts = slide_text_parts(slide)
        low = " | ".join(parts).lower()
        up  = low.upper()
        # Mark content started once we hit a real DAY/DOMAIN section divider
        # Exclude lesson-plan slides that mention days in their body text
        if not content_started:
            is_schedule = 'LESSON PLAN' in up or 'SCHEDULE' in up
            if not is_schedule:
                if re.search(r'\bDAY\s+\d+\b', up) or re.search(r'\bDOMAIN\s+\d+\b', up):
                    if not re.search(r'\bLAB\s+\d+\b', up):
                        content_started = True
        if not content_started:
            continue   # skip admin check for intentional intro slides
        wc = len(re.findall(r'[a-zA-Z]{3,}', low))
        if wc < 2:
            continue
        for p in ADMIN_PATTERNS:
            if p in low:
                flags.append({"check": "3_admin_slide", "slide": i+1,
                              "pattern": p, "text": low[:80]})
                break
    return flags


def check4_exact_md5(slides, prs):
    """No exact pixel duplicate images."""
    seen, flags, remove = {}, [], set()
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if not blob:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in seen:
            flags.append({"check": "4_exact_md5", "slide": i+1,
                          "first_at": seen[h], "action": "remove"})
            remove.add(i)
        else:
            seen[h] = i+1
    return flags, remove


def check5_image_title_dup(slides, prs, pdf_index):
    """Each PDF source title appears at most once (keep first occurrence)."""
    seen_titles, flags, remove = {}, [], set()
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if not blob:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h not in pdf_index:
            continue
        raw = pdf_index[h][2]
        if not raw:
            continue
        n = normalise(raw)
        if not n:
            continue
        if n in seen_titles:
            flags.append({"check": "5_image_title_dup", "slide": i+1,
                          "title": raw, "first_at": seen_titles[n], "action": "remove"})
            remove.add(i)
        else:
            seen_titles[n] = i+1
    return flags, remove


def check6_text_vs_image(slides, prs, pdf_index):
    """Text slides whose topic matches an image slide's PDF title — remove text."""
    image_titles = set()
    for slide in slides:
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if not blob:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in pdf_index:
            raw = pdf_index[h][2]
            if raw:
                image_titles.add(normalise(raw))

    flags, remove = [], set()
    for i, slide in enumerate(slides):
        if is_fullslide_image(slide, prs):
            continue
        all_text = " ".join(slide_text_parts(slide))
        if re.search(r'\bLAB\s+\d+\b', all_text, re.I):
            continue
        text_title = extract_text_slide_title(slide)
        if not text_title:
            continue
        n = normalise(text_title)
        if not n or len(n.split()) < 3:
            continue
        n_words = set(n.split())
        best_score, best_match = 0.0, None
        for m in image_titles:
            m_words = set(m.split())
            if len(m_words) < 3:
                continue
            inter = len(n_words & m_words)
            union = len(n_words | m_words)
            j = inter / union if union else 0
            if j > best_score:
                best_score, best_match = j, m
        if best_score >= 0.70 and best_match:
            flags.append({"check": "6_text_vs_image", "slide": i+1,
                          "text_title": text_title, "image_match": best_match,
                          "jaccard": round(best_score, 2), "action": "remove"})
            remove.add(i)
    return flags, remove


def check7_content_before_labs(slides, prs, pdf_index):
    """For each section, theory images from THAT section must appear BEFORE its first lab header.
    Images from other sections that appear after the current section's first lab are NOT flagged —
    those belong to their own sections and are correctly placed relative to their own first lab.
    """
    # Build slide_idx → set of section names it belongs to, using PDF source mapping
    # pdf_index has: h → (pdf_name, page_1idx, title)  (1-indexed pages)
    slide_sections = {}
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if not blob:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h not in pdf_index:
            continue
        pdf_name, page_1idx, _ = pdf_index[h]
        page_0idx = page_1idx - 1  # convert to 0-indexed to match SECTIONS_PDF
        for sec in SECTIONS_PDF:
            if sec["pdf"] == pdf_name and sec["p_start"] <= page_0idx < sec["p_end"]:
                slide_sections.setdefault(i, set()).add(sec["name"])

    flags = []
    for sec_name, kws in SECTION_ORDER_CHECKS:
        # Find first lab header for this section
        first_lab_idx = None
        for i, slide in enumerate(slides):
            if is_fullslide_image(slide, prs):
                continue
            up = slide_upper(slide)
            if all(kw in up for kw in kws):
                first_lab_idx = i
                break
        if first_lab_idx is None:
            continue

        # Only flag image slides that (a) belong to THIS section AND (b) appear after first_lab_idx
        misplaced = []
        for i in range(first_lab_idx + 1, len(slides)):
            if not is_fullslide_image(slides[i], prs):
                continue
            if sec_name in slide_sections.get(i, set()):
                misplaced.append(i + 1)  # 1-indexed

        if misplaced:
            flags.append({
                "check": "7_content_before_labs",
                "section": sec_name,
                "first_lab_slide": first_lab_idx + 1,
                "misplaced_image_slides": misplaced[:10],
                "total_misplaced": len(misplaced),
                "detail": f"{sec_name}: {len(misplaced)} theory images after first lab (slide {first_lab_idx+1})",
            })
    return flags


# ── main ─────────────────────────────────────────────────────────────────────

def validate(pptx_path, fix_output=None):
    print("Building PDF source index...", file=sys.stderr)
    pdf_index = build_pdf_index()
    print(f"  Indexed {len(pdf_index)} pages\n", file=sys.stderr)

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    total = len(slides)

    all_flags = []
    to_remove = set()

    f1 = check1_lab_completeness(slides, prs)
    f2 = check2_killercoda_urls(slides, prs)
    f3 = check3_admin_slides(slides, prs)
    f4, r4 = check4_exact_md5(slides, prs)
    f5, r5 = check5_image_title_dup(slides, prs, pdf_index)
    f6, r6 = check6_text_vs_image(slides, prs, pdf_index)
    f7 = check7_content_before_labs(slides, prs, pdf_index)

    all_flags.extend(f1 + f2 + f3 + f4 + f5 + f6 + f7)
    to_remove = r4 | r5 | r6

    print(f"Check 1 — Lab completeness:      {'PASS' if not f1 else f'FAIL ({len(f1)} issues)'}", file=sys.stderr)
    print(f"Check 2 — KillerCoda URLs:        {'PASS' if not f2 else f'FAIL ({len(f2)} missing)'}", file=sys.stderr)
    print(f"Check 3 — Admin slides:           {'PASS' if not f3 else f'FAIL ({len(f3)} found)'}", file=sys.stderr)
    print(f"Check 4 — Exact MD5 duplicates:   {'PASS' if not f4 else f'FAIL ({len(f4)} dups)'}", file=sys.stderr)
    print(f"Check 5 — Image title duplicates: {'PASS' if not f5 else f'FAIL ({len(f5)} dups)'}", file=sys.stderr)
    print(f"Check 6 — Text-vs-image overlap:  {'PASS' if not f6 else f'FAIL ({len(f6)} text slides)'}", file=sys.stderr)
    print(f"Check 7 — Content before labs:    {'PASS' if not f7 else f'FAIL ({len(f7)} sections out of order)'}", file=sys.stderr)

    if fix_output and to_remove:
        remove_slides(prs, to_remove)
        prs.save(fix_output)
        total_out = len(list(prs.slides))
        print(f"\n  Auto-fix applied: removed {len(to_remove)} slides → {total_out} remaining", file=sys.stderr)
        print(f"  Saved to: {fix_output}", file=sys.stderr)
    elif fix_output:
        import shutil
        shutil.copy(pptx_path, fix_output)
        total_out = total
    else:
        total_out = total

    passes = sum(1 for f in [f1,f2,f3,f4,f5,f6,f7] if not f)
    result = {
        "file": pptx_path,
        "total_slides_in": total,
        "total_slides_out": total_out,
        "checks_passed": passes,
        "checks_failed": 7 - passes,
        "all_pass": len(all_flags) == 0,
        "flags": all_flags,
    }
    return result


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input", help="Input PPTX")
    parser.add_argument("--fix-output", help="If given, auto-fix checks 4-6 and save here")
    args = parser.parse_args()

    result = validate(args.input, fix_output=args.fix_output)
    print(json.dumps(result, indent=2))

    status = "ALL PASS" if result["all_pass"] else f"{result['checks_failed']} CHECK(S) FAILED"
    print(f"\nVALIDATION {status} — {result['total_slides_in']} slides in, "
          f"{result['checks_passed']}/7 checks passed", file=sys.stderr)
    sys.exit(0 if result["all_pass"] else 1)
