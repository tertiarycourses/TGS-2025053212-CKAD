"""
CKAD Slide Reorder
==================
Moves PDF-image theory slides to appear BEFORE the first lab header of their
section.  The build script (build_ckad_v4.py) currently injects images AFTER
the last TEST-IT slide, producing: [Lab slides] → [Theory images].
This script corrects that to: [Theory images] → [Lab slides].

Usage:
    python reorder_ckad.py <input.pptx> [output.pptx]

If output is omitted the input is overwritten.
"""
import sys, os, hashlib, json, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"

# PDF source files and their page ranges per section
# (pdf_key, page_start 0-indexed inclusive, page_end exclusive)
SOURCE_PDFS = {
    "Day1":  BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf",
    "Day2":  BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf",
    "Day3":  BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf",
    "Day4M": BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf",
    "Day4A": BASE_DIR + r"\CKAD Day 4 Afternoon - Summary & Tips.pdf",
}

# Section definitions: name, pdf key, page range, and keywords to find the first lab header
# first_lab: list of strings that ALL appear (uppercase) in the first lab header slide for this section
SECTIONS = [
    {"name": "Container Images",       "pdf": "Day1",  "p_start": 0,   "p_end": 28,  "first_lab": ["CONTAINER IMAGES", "LAB 1"]},
    {"name": "Pods",                   "pdf": "Day1",  "p_start": 28,  "p_end": 82,  "first_lab": ["PODS", "LAB 3"]},
    {"name": "Jobs",                   "pdf": "Day1",  "p_start": 82,  "p_end": 97,  "first_lab": ["JOBS", "LAB 4"]},
    {"name": "CronJobs",               "pdf": "Day1",  "p_start": 97,  "p_end": 106, "first_lab": ["CRONJOBS", "LAB 5"]},
    {"name": "Multi-container Pods",   "pdf": "Day3",  "p_start": 98,  "p_end": 113, "first_lab": ["MULTI-CONTAINER PODS", "LAB 6"]},
    {"name": "Volumes",                "pdf": "Day3",  "p_start": 10,  "p_end": 33,  "first_lab": ["VOLUMES", "LAB 8"]},
    {"name": "Labels and Annotations", "pdf": "Day1",  "p_start": 106, "p_end": 120, "first_lab": None},
    {"name": "Deployments",            "pdf": "Day1",  "p_start": 122, "p_end": 136, "first_lab": ["DEPLOYMENTS", "LAB 9"]},
    {"name": "Rolling Updates",        "pdf": "Day1",  "p_start": 136, "p_end": 155, "first_lab": ["ROLLING UPDATES", "LAB 10"]},
    {"name": "Blue/Green and Canary",  "pdf": "Day1",  "p_start": 157, "p_end": 164, "first_lab": ["BLUE/GREEN", "LAB 11"]},
    {"name": "HPA (Autoscaler)",       "pdf": "Day1",  "p_start": 164, "p_end": 170, "first_lab": None},
    {"name": "Helm",                   "pdf": "Day3",  "p_start": 34,  "p_end": 48,  "first_lab": ["HELM", "LAB 13"]},
    {"name": "Kustomize",              "pdf": "Day3",  "p_start": 48,  "p_end": 62,  "first_lab": ["KUSTOMIZE", "LAB 14"]},
    {"name": "DaemonSets/StatefulSets","pdf": "Day4M", "p_start": 9,   "p_end": 20,  "first_lab": ["DAEMON", "LAB 15"]},
    {"name": "Probes",                 "pdf": "Day3",  "p_start": 80,  "p_end": 92,  "first_lab": ["PROBES", "LAB 16"]},
    {"name": "Metrics Server",         "pdf": "Day3",  "p_start": 92,  "p_end": 98,  "first_lab": ["METRICS", "LAB 18"]},
    {"name": "API Deprecations",       "pdf": "Day3",  "p_start": 72,  "p_end": 80,  "first_lab": ["API VERSIONS", "LAB 20"]},
    {"name": "ConfigMaps",             "pdf": "Day2",  "p_start": 130, "p_end": 141, "first_lab": ["CONFIGMAPS", "LAB 21"]},
    {"name": "Secrets",                "pdf": "Day2",  "p_start": 141, "p_end": 152, "first_lab": ["SECRETS", "LAB 22"]},
    {"name": "SecurityContext",        "pdf": "Day2",  "p_start": 153, "p_end": 161, "first_lab": ["SECURITYCONTEXT", "LAB 23"]},
    {"name": "ServiceAccounts",        "pdf": "Day2",  "p_start": 100, "p_end": 107, "first_lab": ["SERVICEACCOUNTS", "LAB 24"]},
    {"name": "RBAC",                   "pdf": "Day2",  "p_start": 83,  "p_end": 100, "first_lab": ["RBAC", "LAB 25"]},
    {"name": "CRDs",                   "pdf": "Day3",  "p_start": 62,  "p_end": 72,  "first_lab": None},
    {"name": "Quotas and Limits",      "pdf": "Day2",  "p_start": 107, "p_end": 129, "first_lab": ["QUOTAS", "LAB 26"]},
    {"name": "Taints and Node Affinity","pdf": "Day4M","p_start": 20,  "p_end": 34,  "first_lab": None},
    {"name": "Services",               "pdf": "Day2",  "p_start": 22,  "p_end": 55,  "first_lab": ["SERVICES", "LAB 27"]},
    {"name": "Ingress and TLS",        "pdf": "Day2",  "p_start": 55,  "p_end": 67,  "first_lab": ["INGRESS", "LAB 29"]},
    {"name": "NetworkPolicy",          "pdf": "Day2",  "p_start": 67,  "p_end": 82,  "first_lab": ["NETWORKPOLICY", "LAB 30"]},
]
RENDER_DPI = 2.0


# ── helpers ─────────────────────────────────────────────────────────────────

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def get_img_blob(slide):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh.image.blob
    return None

def slide_upper_text(slide):
    return " ".join(
        sh.text_frame.text.upper()
        for sh in slide.shapes if sh.has_text_frame
    )

def build_pdf_index():
    """MD5(rendered page) → (pdf_key, page_num_0indexed, title)."""
    index = {}
    for pdf_key, path in SOURCE_PDFS.items():
        if not os.path.exists(path):
            print(f"  [WARN] PDF not found: {path}", file=sys.stderr)
            continue
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(RENDER_DPI, RENDER_DPI)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            h = hashlib.md5(pix.tobytes("png")).hexdigest()
            lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
            title = lines[0][:60] if lines else ""
            index[h] = (pdf_key, i, title)   # i is 0-indexed
        doc.close()
    return index


def find_first_slide_idx(slides, keywords):
    """Return the 0-based index of the first slide whose upper text contains ALL keywords."""
    for i, slide in enumerate(slides):
        up = slide_upper_text(slide)
        if all(kw in up for kw in keywords):
            return i
    return None


def reorder_slides(prs, new_order):
    """Apply new_order (list of 0-based old indices) to prs slides."""
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    ordered = [all_ids[i] for i in new_order]
    for el in list(sp):
        sp.remove(el)
    for el in ordered:
        sp.append(el)


# ── main ─────────────────────────────────────────────────────────────────────

def reorder(input_path, output_path=None):
    if output_path is None:
        output_path = input_path

    print("Building PDF source index...", file=sys.stderr)
    pdf_index = build_pdf_index()
    print(f"  Indexed {len(pdf_index)} PDF pages", file=sys.stderr)

    prs = Presentation(input_path)
    slides = list(prs.slides)
    n = len(slides)

    # Map each image slide index → (pdf_key, page_0idx)
    img_source = {}
    for i, slide in enumerate(slides):
        if is_fullslide_image(slide, prs):
            blob = get_img_blob(slide)
            if blob:
                h = hashlib.md5(blob).hexdigest()
                if h in pdf_index:
                    img_source[i] = pdf_index[h]   # (pdf_key, page_0idx, title)

    # Build set of slide indices belonging to each section
    section_images = {}  # section_name → sorted list of slide indices
    for sec in SECTIONS:
        name = sec["name"]
        pdf_key = sec["pdf"]
        p_start = sec["p_start"]
        p_end = sec["p_end"]
        members = [
            i for i, (pk, pg, _) in img_source.items()
            if pk == pdf_key and p_start <= pg < p_end
        ]
        section_images[name] = sorted(members)

    # For each section, find its first lab slide (target insertion point)
    section_first_lab = {}
    for sec in SECTIONS:
        name = sec["name"]
        kws = sec["first_lab"]
        if kws is None:
            section_first_lab[name] = None
            continue
        idx = find_first_slide_idx(slides, kws)
        section_first_lab[name] = idx

    # Determine which image slides need to move and where
    moves = []  # (current_idx, insert_before_idx)
    for sec in SECTIONS:
        name = sec["name"]
        first_lab = section_first_lab[name]
        if first_lab is None:
            continue
        img_slides = section_images.get(name, [])
        # Image slides that currently appear AFTER the first lab header need to move
        misplaced = [i for i in img_slides if i > first_lab]
        if misplaced:
            print(f"  [{name}] Moving {len(misplaced)} images to before slide {first_lab+1}", file=sys.stderr)
            for mi in misplaced:
                moves.append((mi, first_lab))

    if not moves:
        print("  No ordering issues found — deck already correct.", file=sys.stderr)
        if output_path != input_path:
            import shutil
            shutil.copy(input_path, output_path)
        return {"moves": 0, "output": output_path}

    # Build the new slide order by applying all moves simultaneously
    # For each section, collect its misplaced images, move them before first_lab
    # Strategy: build new_order list by iterating through slides and inserting groups
    misplaced_set = {mi for mi, _ in moves}

    # Group moves by target (insert_before_idx)
    from collections import defaultdict
    groups = defaultdict(list)
    for mi, before in moves:
        groups[before].append(mi)

    new_order = []
    for i in range(n):
        if i in misplaced_set:
            continue  # will be inserted at the target position
        # Before inserting slide i, check if any group targets i
        if i in groups:
            # Insert the misplaced images from this group (in their original relative order)
            group_sorted = sorted(groups[i])
            new_order.extend(group_sorted)
        new_order.append(i)

    assert len(new_order) == n, f"Order mismatch: {len(new_order)} vs {n}"

    reorder_slides(prs, new_order)
    prs.save(output_path)

    total_moved = len(misplaced_set)
    print(f"\n  Reordered {total_moved} slides.", file=sys.stderr)
    return {"moves": total_moved, "output": output_path}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python reorder_ckad.py <input.pptx> [output.pptx]")
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else inp
    result = reorder(inp, out)
    print(json.dumps(result, indent=2))
