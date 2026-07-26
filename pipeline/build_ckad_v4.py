"""
CKAD v4.0 Full Build Pipeline
==============================
Steps:
  1. Strip all full-slide PDF images from v3.3 → clean designed base
  2. Domain 1  — import images for Labs 1-8   (Container Images, Pods, Jobs,
                  CronJobs, Multi-container, Volumes) from Day1 + Day3 PDFs
  3. Domain 2  — import images for Labs 9-15  (Labels, Deployments, Rolling,
                  Blue/Green, Canary, HPA, Helm, Kustomize, DaemonSets)
  4. Domain 3  — import images for Labs 16-20 (Probes, Logging, Metrics,
                  Debugging, API Deprecations)
  5. Domain 4  — import images for Labs 21-26 (ConfigMaps, Secrets,
                  SecurityContext, ServiceAccounts, RBAC, Quotas)
  6. Domain 5  — import images for Labs 27-30 + Taints/Node Affinity
                  (Services, Ingress, NetworkPolicy) from Day2 + Day4M PDFs
  7. Audit after every domain — auto-fix and loop until PASS

Run:
    python pipeline/build_ckad_v4.py
"""
import sys, os, re, io, json
sys.stdout.reconfigure(encoding='utf-8')

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SRC_V33  = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v3.3.pptx"
OUT_V4   = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v4.0.pptx"
TMP_DIR  = BASE_DIR + r"\pipeline\_tmp"
os.makedirs(TMP_DIR, exist_ok=True)

PDF1  = BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf"
PDF2  = BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf"
PDF3  = BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf"
PDF4M = BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf"

RENDER_DPI = 2.0

SKIP_PATTERNS = [
    "exam essentials","exercise","· commands","see you tomorrow",
    "thank you everyone","basic ground rules","admin","topics we",
    "tea break","lunch break","wrap-up","your trainer","traqom",
    "attendance","topic 1","topic 2","topic 3","topic 4","topic 5","topic 6",
    "topic 7","topic 8","topic 9","topic 10","topic 11",
    "lab 1 ·","lab 2 ·","lab 3 ·","lab 4 ·","lab 5 ·","lab 6 ·","lab 7 ·",
    "lab 8 ·","lab 9 ·","lab 10 ·","lab 11 ·","lab 12 ·","lab 13 ·",
    "lab 14 ·","lab 15 ·","lab 16 ·","lab 17 ·","lab 18 ·","lab 19 ·",
    "lab 20 ·","lab 21 ·","lab 22 ·","lab 23 ·","lab 24 ·","lab 25 ·",
    "lab 26 ·","lab 27 ·","lab 28 ·","lab 29 ·","lab 30 ·",
    "assessment rules","final assessment","course slides","companion study guide",
    "video course","preparing for the exam","key takeaways","summary","recap",
    "authentication and authorization",
]

# Domain sections: name, pdf_path, page_start(0-idx), page_end(excl), find_keywords
DOMAIN_SECTIONS = [
    # ── Domain 1 ──────────────────────────────────────────────────────────────
    {"name": "Container Images",       "pdf": PDF1, "p_start": 0,   "p_end": 28,
     "find": ["MULTI-STAGE", "TEST IT"]},
    {"name": "Pods",                   "pdf": PDF1, "p_start": 28,  "p_end": 82,
     "find": ["MANAGE PODS", "TEST IT"]},
    {"name": "Jobs",                   "pdf": PDF1, "p_start": 82,  "p_end": 97,
     "find": ["RUN-TO-COMPLETION", "TEST IT"]},
    {"name": "CronJobs",               "pdf": PDF1, "p_start": 97,  "p_end": 106,
     "find": ["SCHEDULED WORKLOADS", "TEST IT"]},
    {"name": "Multi-container Pods",   "pdf": PDF3, "p_start": 98,  "p_end": 113,
     "find": ["INIT CONTAINERS", "TEST IT"]},
    {"name": "Volumes",                "pdf": PDF3, "p_start": 10,  "p_end": 33,
     "find": ["EMPTYDIR", "TEST IT"]},
    # ── Domain 2 ──────────────────────────────────────────────────────────────
    {"name": "Labels and Annotations", "pdf": PDF1, "p_start": 106, "p_end": 120,
     "find": ["KEY TAKEAWAYS", "LABELS ARE KEY-VALUE"]},
    {"name": "Deployments",            "pdf": PDF1, "p_start": 122, "p_end": 136,
     "find": ["DEPLOYMENTS AND REPLICASETS", "TEST IT"]},
    {"name": "Rolling Updates",        "pdf": PDF1, "p_start": 136, "p_end": 155,
     "find": ["ROLLING UPDATES AND ROLLBACK", "TEST IT"]},
    {"name": "Blue/Green and Canary",  "pdf": PDF1, "p_start": 157, "p_end": 164,
     "find": ["CANARY DEPLOYMENT", "TEST IT"]},
    {"name": "HPA (Autoscaler)",       "pdf": PDF1, "p_start": 164, "p_end": 170,
     "find": ["HORIZONTAL POD AUTOSCALER", "EXERCISE"]},
    {"name": "Helm",                   "pdf": PDF3, "p_start": 34,  "p_end": 48,
     "find": ["INSTALL, UPGRADE, ROLLBACK", "TEST IT"]},
    {"name": "Kustomize",              "pdf": PDF3, "p_start": 48,  "p_end": 62,
     "find": ["KUSTOMIZE OVERLAYS", "TEST IT"]},
    {"name": "DaemonSets/StatefulSets","pdf": PDF4M,"p_start": 9,   "p_end": 20,
     "find": ["DAEMONSETS AND STATEFULSETS", "TEST IT"]},
    # ── Domain 3 ──────────────────────────────────────────────────────────────
    {"name": "Probes",                 "pdf": PDF3, "p_start": 80,  "p_end": 92,
     "find": ["LIVENESS, READINESS & STARTUP PROBES", "TEST IT"]},
    {"name": "Metrics Server",         "pdf": PDF3, "p_start": 92,  "p_end": 98,
     "find": ["KUBECTL TOP AND METRICS SERVER", "TEST IT"]},
    {"name": "API Deprecations",       "pdf": PDF3, "p_start": 72,  "p_end": 80,
     "find": ["API DEPRECATIONS", "TEST IT"]},
    # ── Domain 4 ──────────────────────────────────────────────────────────────
    {"name": "ConfigMaps",             "pdf": PDF2, "p_start": 130, "p_end": 141,
     "find": ["ENV & VOLUME INJECTION", "TEST IT"]},
    {"name": "Secrets",                "pdf": PDF2, "p_start": 141, "p_end": 152,
     "find": ["DB_PASS", "TEST IT"]},
    {"name": "SecurityContext",        "pdf": PDF2, "p_start": 153, "p_end": 161,
     "find": ["SECURITYCONTEXT", "TEST IT"]},
    {"name": "ServiceAccounts",        "pdf": PDF2, "p_start": 100, "p_end": 107,
     "find": ["SERVICEACCOUNTS", "TEST IT"]},
    {"name": "RBAC",                   "pdf": PDF2, "p_start": 83,  "p_end": 100,
     "find": ["ROLES & ROLEBINDINGS", "TEST IT"]},
    {"name": "CRDs",                   "pdf": PDF3, "p_start": 62,  "p_end": 72,
     "find": ["KEY TAKEAWAYS", "IMPLEMENT A CRD"]},
    {"name": "Quotas and Limits",      "pdf": PDF2, "p_start": 107, "p_end": 129,
     "find": ["RESOURCEQUOTA & LIMITRANGE", "TEST IT"]},
    # ── Domain 5 ──────────────────────────────────────────────────────────────
    {"name": "Taints and Node Affinity","pdf": PDF4M,"p_start": 20, "p_end": 34,
     "find": ["DAEMONSETS AND STATEFULSETS", "TEST IT"]},
    {"name": "Services",               "pdf": PDF2, "p_start": 22,  "p_end": 55,
     "find": ["CLUSTERIP, NODEPORT, LOADBALANCER", "TEST IT"]},
    {"name": "Ingress and TLS",        "pdf": PDF2, "p_start": 55,  "p_end": 67,
     "find": ["INGRESS WITH TLS", "TEST IT"]},
    {"name": "NetworkPolicy",          "pdf": PDF2, "p_start": 67,  "p_end": 82,
     "find": ["NETWORKPOLICY", "TEST IT"]},
]


# ── Helpers ──────────────────────────────────────────────────────────────────

def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def is_skip(t):
    return any(p in t for p in SKIP_PATTERNS)

def find_last_slide(slides, search_list):
    last = None
    for i, slide in enumerate(slides):
        combined = " ".join(
            sh.text_frame.text.upper()
            for sh in slide.shapes if sh.has_text_frame
        )
        if all(s in combined for s in search_list):
            last = i
    return last

def collect_pages(pdf_path, p_start, p_end, name):
    doc  = fitz.open(pdf_path)
    keep = []
    src  = os.path.basename(pdf_path)[:20]
    print(f"    [{name}] {src} p{p_start+1}–{p_end}")
    for i in range(p_start, min(p_end, len(doc))):
        page  = doc[i]
        text  = page.get_text()
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        title = lines[0][:60] if lines else "(empty)"
        wc    = len(re.findall(r'[a-zA-Z]{4,}', text.lower()))
        t     = title.lower()
        if wc < 2:
            print(f"      p{i+1:3d} BLANK")
            continue
        if is_skip(t):
            print(f"      p{i+1:3d} SKIP  {title}")
            continue
        print(f"      p{i+1:3d} KEEP  {title}")
        keep.append(i)
    doc.close()
    return keep

def render_and_inject(prs, pdf_path, pages, inject_after):
    SW, SH = prs.slide_width, prs.slide_height
    old_n  = len(prs.slides)

    for pi in pages:
        doc = fitz.open(pdf_path)
        page = doc[pi]
        mat  = fitz.Matrix(RENDER_DPI, RENDER_DPI)
        pix  = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png")
        doc.close()

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic   = slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, SW, SH)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        footer_h = int(SH * 0.07)
        box = slide.shapes.add_shape(1, 0, SH - footer_h, SW, footer_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.fill.background()

    sp     = prs.slides._sldIdLst
    all_id = list(sp)
    exist  = all_id[:old_n]
    new_id = all_id[old_n:]
    ordered = exist[:inject_after+1] + new_id + exist[inject_after+1:]
    for el in list(sp): sp.remove(el)
    for el in ordered:  sp.append(el)

def run_audit(pptx_path):
    # Inline audit (same rules as audit_ckad.py)
    from pipeline.audit_ckad import audit
    return audit(pptx_path)


# ── Step 1: Strip images from v3.3 ───────────────────────────────────────────

print(f"\n{'='*65}")
print("STEP 1 — Strip PDF images from v3.3 → designed base")
print(f"{'='*65}")

prs = Presentation(SRC_V33)
slides = list(prs.slides)
img_count = sum(1 for s in slides if is_fullslide_image(s, prs))
print(f"  v3.3: {len(slides)} slides, {img_count} PDF-image slides to strip")

remove_idx = {i for i, s in enumerate(slides) if is_fullslide_image(s, prs)}
sp = prs.slides._sldIdLst
all_ids = list(sp)
kept = [el for j, el in enumerate(all_ids) if j not in remove_idx]
for el in list(sp): sp.remove(el)
for el in kept: sp.append(el)

BASE_PATH = TMP_DIR + r"\v4_base.pptx"
prs.save(BASE_PATH)
print(f"  Base saved: {len(kept)} designed slides → v4_base.pptx")


# ── Steps 2-6: Import domain images ──────────────────────────────────────────

print(f"\n{'='*65}")
print("STEPS 2-6 — Import PDF images domain by domain")
print(f"{'='*65}")

prs = Presentation(BASE_PATH)
total_added = 0

for sec in DOMAIN_SECTIONS:
    name    = sec["name"]
    find_kw = sec["find"]

    inject_after = find_last_slide(prs.slides, find_kw)
    if inject_after is None:
        print(f"\n  WARNING: {find_kw} not found — skipping '{name}'")
        continue

    pages = collect_pages(sec["pdf"], sec["p_start"], sec["p_end"], name)
    if not pages:
        print(f"    No pages to inject for '{name}'")
        continue

    before = len(prs.slides)
    render_and_inject(prs, sec["pdf"], pages, inject_after)
    added = len(prs.slides) - before
    total_added += added
    print(f"    → Injected {added} slides after slide {inject_after+1}. Total: {len(prs.slides)}")

DRAFT_PATH = TMP_DIR + r"\v4_draft.pptx"
print(f"\n  Saving draft → v4_draft.pptx  ({len(list(prs.slides))} slides)")
prs.save(DRAFT_PATH)
print(f"  Total images added: {total_added}")


# ── Step 6b: Deduplicate image slides ────────────────────────────────────────

print(f"\n{'='*65}")
print("STEP 6b — Deduplicate identical image slides")
print(f"{'='*65}")

DEDUP_SCRIPT = os.path.join(os.path.dirname(__file__), "dedupe_ckad.py")
DEDUPED_PATH = TMP_DIR + r"\v4_deduped.pptx"

dedup_result = subprocess.run(
    [PY, DEDUP_SCRIPT, DRAFT_PATH, DEDUPED_PATH],
    capture_output=True, text=True, encoding='utf-8', errors='replace'
)
try:
    dedup_report = json.loads(dedup_result.stdout)
    removed = dedup_report["duplicates_removed"]
    print(f"  Duplicates removed: {removed}")
    if removed:
        print(f"  Removed slides: {dedup_report['removed_slides']}")
    print(f"  Slides after dedupe: {dedup_report['total_out']}")
    current_for_audit = DEDUPED_PATH
except json.JSONDecodeError:
    print("  [DEDUPE WARNING] Could not parse output — using draft as-is")
    current_for_audit = DRAFT_PATH


# ── Step 7: Audit + fix loop ─────────────────────────────────────────────────

print(f"\n{'='*65}")
print("STEP 7 — Audit + fix loop")
print(f"{'='*65}")

import subprocess
AUDIT_SCRIPT = os.path.join(os.path.dirname(__file__), "audit_ckad.py")
FIX_SCRIPT   = os.path.join(os.path.dirname(__file__), "fix_ckad.py")
PY = sys.executable

current = current_for_audit
for iteration in range(1, 6):
    print(f"\n  ── Iteration {iteration} ──")
    result = subprocess.run([PY, AUDIT_SCRIPT, current],
                            capture_output=True, text=True, encoding='utf-8', errors='replace')
    try:
        report = json.loads(result.stdout)
    except json.JSONDecodeError:
        print("  [AUDIT ERROR] could not parse JSON — stopping")
        break

    n_viol = len(report["violations"])
    print(f"  Slides: {report['total_slides']}  Labs: {len(report['labs_found'])}  Violations: {n_viol}")

    if report["pass"]:
        print("  ✓ AUDIT PASS")
        break

    for v in report["violations"]:
        print(f"    • [{v['rule']}] slide {v['slide']}: {v.get('detail','')[:70]}")

    rpt_path = TMP_DIR + f"\\audit_iter{iteration}.json"
    with open(rpt_path, 'w', encoding='utf-8') as f:
        json.dump(report, f, indent=2)

    fixed_path = TMP_DIR + f"\\v4_fixed_iter{iteration}.pptx"
    subprocess.run([PY, FIX_SCRIPT, current, rpt_path, fixed_path])
    current = fixed_path
else:
    print("  ✗ Still failing after 5 iterations")

# ── Final save ───────────────────────────────────────────────────────────────
import shutil
shutil.copy(current, OUT_V4)

print(f"\n{'='*65}")
print(f"OUTPUT → {os.path.basename(OUT_V4)}")
final_prs = Presentation(OUT_V4)
print(f"FINAL:  {len(list(final_prs.slides))} slides")
print(f"{'='*65}\n")
