"""
replace_admin_slides.py
========================
Replaces the content of slides 0-5 in v4.5-final in-place with fresh
content from v2.0.  Uses get_or_add_image_part so icons display correctly.
No text-colour changes — preserves v2.0 dark-on-white style.

This avoids the delete+re-add approach which causes ZIP filename conflicts.

Usage:
    python pipeline/replace_admin_slides.py
"""
import sys, copy
from io import BytesIO
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SRC_V20  = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v2.0.pptx"
IN_V45   = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v4.5-final.pptx"
OUT_V45  = BASE_DIR + r"\courseware\CKAD-Certified-Kubernetes-Application-Developer-v4.5-admin2.pptx"

# Map: destination slot (0-based) → v2.0 source index (0-based)
ADMIN_MAP = {
    0: 1,   # Welcome & Housekeeping
    1: 2,   # Digital Attendance
    2: 3,   # About the Trainer
    3: 7,   # Lesson Plan / Schedule
    4: 9,   # Final Assessment
    5: 10,  # Assessment Rules
}

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def replace_slide_content(dst_slide, src_slide):
    """
    Replace dst_slide's shape content with src_slide's shapes.
    Properly transfers image parts so icons display correctly.
    Text colours are taken as-is from src_slide (dark-on-white).
    """
    # Transfer image parts from src_slide to dst_slide
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.is_external or "/image" not in rel.reltype:
            continue
        try:
            blob = rel.target_part.blob
            _img_part, new_rId = dst_slide.part.get_or_add_image_part(BytesIO(blob))
            rId_map[rId] = new_rId
        except Exception as e:
            print(f"      [WARN] image rId {rId}: {e}")

    # Deep-copy src spTree, rewriting rIds, then replace dst spTree content
    dst_sp = dst_slide.shapes._spTree
    src_sp = src_slide.shapes._spTree

    for ch in list(dst_sp):
        dst_sp.remove(ch)

    for ch in src_sp:
        new_ch = copy.deepcopy(ch)
        for el in new_ch.iter():
            for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
                v = el.get(attr)
                if v and v in rId_map:
                    el.set(attr, rId_map[v])
        dst_sp.append(new_ch)


def slide_label(prs, idx):
    parts = []
    for sh in prs.slides[idx].shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()[:50]
            if t:
                parts.append(t)
    return " | ".join(parts[:2])


def main():
    print("Opening presentations ...")
    src = Presentation(SRC_V20)
    dst = Presentation(IN_V45)
    print(f"  v2.0 : {len(list(src.slides))} slides")
    print(f"  v4.5 : {len(list(dst.slides))} slides")

    print(f"\nReplacing first {len(ADMIN_MAP)} slides in-place ...")
    for dst_idx, src_idx in sorted(ADMIN_MAP.items()):
        src_lbl = slide_label(src, src_idx)
        dst_lbl = slide_label(dst, dst_idx)
        print(f"  Slot {dst_idx+1}: [{dst_lbl[:40]}] <- v2.0[{src_idx+1}] [{src_lbl[:40]}]")
        replace_slide_content(dst.slides[dst_idx], src.slides[src_idx])

    dst.save(OUT_V45)
    print(f"\nSaved -> {OUT_V45}  ({len(list(dst.slides))} slides total)")


if __name__ == "__main__":
    main()
