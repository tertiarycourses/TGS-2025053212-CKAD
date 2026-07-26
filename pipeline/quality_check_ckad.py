"""
CKAD Slide Quality Check Agent
================================
Comprehensive duplicate / redundancy detector.  Three checks:

  Check A — TEXT vs IMAGE same-topic:
    For each designed text slide, extract its topic title and see if an
    IMAGE slide covers the same topic.  Flag the TEXT slide for removal
    (prefer visual PDF slide over text-only).

  Check B — IMAGE vs IMAGE cross-section duplicates:
    For image slides whose source title appears more than once in the deck,
    keep the first occurrence, flag later ones for removal.
    Uses the full PDF source index (all 5 PDFs including Day4 Afternoon).

  Check C — Exact pixel duplicates (MD5):
    Belt-and-suspenders pass — catches any exact copies missed earlier.

Usage:
    python quality_check_ckad.py <deck.pptx> [output.pptx] [--dry-run]

  --dry-run : Report only, do not save.  Prints a full report to stdout.

Loops until zero removals (convergence).  Always exits 0.
"""
import sys, os, hashlib, json, re, argparse
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

BASE_DIR = r"C:\Users\mohan\agents\CKAD-2026"
SOURCE_PDFS = {
    "Day1":  BASE_DIR + r"\CKAD Day 1 - Intro, pods, ns, jobs, cronjobs, labels, deployments.pdf",
    "Day2":  BASE_DIR + r"\CKAD Day 2 - Services, RBAC, Resource Limits, Configmaps & Secrets.pdf",
    "Day3":  BASE_DIR + r"\CKAD Day 3 - Volumes, Probes, CRD, Multi-Container Pods.pdf",
    "Day4M": BASE_DIR + r"\CKAD Day 4 Morning - DaemonSets, StatefulSets.pdf",
    "Day4A": BASE_DIR + r"\CKAD Day 4 Afternoon - Summary & Tips.pdf",
}
RENDER_DPI = 2.0

# Prefixes stripped from designed text slide titles before matching
TEXT_TITLE_NOISE = [
    r"^(DOMAIN\s+\d+\s*[·:—\-]+\s*)",  # "DOMAIN 1 · "
    r"^(DAY\s+\d+\s*[·:—\-]+\s*)",      # "DAY 1 · "
    r"^(FOUNDATIONS\s*[|·:—\-]+\s*)",    # "FOUNDATIONS | "
    r"^(THEORY\s*[|·:—\-]+\s*)",
    r"^(CONCEPTS?\s*[|·:—\-]+\s*)",
    r"^(OVERVIEW\s*[|·:—\-]+\s*)",
    r"^(SECTION\s+\d+\s*[|·:—\-]+\s*)",
]

# ── helpers ─────────────────────────────────────────────────────────────────

def normalise(s):
    """Lowercase, strip punctuation runs, collapse spaces."""
    s = s.lower().strip()
    for pat in TEXT_TITLE_NOISE:
        s = re.sub(pat, '', s, flags=re.IGNORECASE)
    s = re.sub(r'[^a-z0-9 ]', ' ', s)
    return re.sub(r'\s+', ' ', s).strip()


def build_pdf_index():
    """MD5(rendered page bytes) → (pdf_name, page_num, raw_title)."""
    index = {}
    for pdf_name, path in SOURCE_PDFS.items():
        if not os.path.exists(path):
            print(f"  [WARN] PDF not found: {os.path.basename(path)}", file=sys.stderr)
            continue
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(RENDER_DPI, RENDER_DPI)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            blob = pix.tobytes("png")
            h = hashlib.md5(blob).hexdigest()
            lines = [l.strip() for l in page.get_text().split('\n') if l.strip()]
            title = lines[0][:80] if lines else ""
            index[h] = (pdf_name, i + 1, title)
        doc.close()
    return index


def is_fullslide_image(slide, prs):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.width > int(prs.slide_width * 0.85):
                return True
    return False

def get_img_blob(slide):
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh.image.blob
    return None

def slide_text_parts(slide):
    """Return list of non-empty text strings from all text frames."""
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                parts.append(t)
    return parts

def extract_text_slide_title(slide):
    """
    For a designed text slide, find the topic title (not the eyebrow category label).
    Collects all candidate lines, then returns the first one that has >= 3 words
    after normalising — this skips single-word eyebrow labels like "FOUNDATIONS".
    """
    parts = slide_text_parts(slide)
    if not parts:
        return None

    candidates = []
    for part in parts:
        lines = [l.strip() for l in part.split('\n') if l.strip()]
        for line in lines:
            if len(line) < 8:
                continue
            if re.match(r'^(LAB\s+\d+|STEP\s+\d+|DAY\s+\d+\s*·)', line, re.I):
                continue
            if re.match(r'^(Tertiary|©|http)', line, re.I):
                continue
            candidates.append(line)

    # Prefer a candidate whose normalised form has >= 3 words (skips one-word labels)
    for c in candidates:
        if len(normalise(c).split()) >= 3:
            return c
    return candidates[0] if candidates else None


def remove_slides(prs, remove_idx):
    sp = prs.slides._sldIdLst
    all_ids = list(sp)
    kept = [el for j, el in enumerate(all_ids) if j not in remove_idx]
    for el in list(sp):
        sp.remove(el)
    for el in kept:
        sp.append(el)


# ── checks ───────────────────────────────────────────────────────────────────

def check_a_text_vs_image(prs, pdf_index):
    """
    Find TEXT slides whose topic title matches an IMAGE slide's PDF title.
    Flag the TEXT slide for removal (keep the visual).
    """
    slides = list(prs.slides)
    remove = set()
    flags = []

    # Build set of all image slide titles (normalised)
    image_titles = set()
    for slide in slides:
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in pdf_index:
            raw = pdf_index[h][2]
            if raw:
                image_titles.add(normalise(raw))

    # Scan text slides for matching titles
    for i, slide in enumerate(slides):
        if is_fullslide_image(slide, prs):
            continue  # skip image slides
        # Never remove lab header / lab command slides
        all_text = " ".join(slide_text_parts(slide))
        if re.search(r'\bLAB\s+\d+\b', all_text, re.I):
            continue
        text_title = extract_text_slide_title(slide)
        if not text_title:
            continue
        n = normalise(text_title)
        if not n:
            continue
        # Check if any image slide has the same topic (Jaccard >= 0.70)
        n_words = set(n.split())
        if len(n_words) < 3:
            continue  # too short to match reliably
        best_score, best_match = 0.0, None
        for m in image_titles:
            m_words = set(m.split())
            if len(m_words) < 3:
                continue
            inter = len(n_words & m_words)
            union = len(n_words | m_words)
            jaccard = inter / union if union else 0
            if jaccard > best_score:
                best_score, best_match = jaccard, m
        if best_score >= 0.70 and best_match:
            remove.add(i)
            flags.append({
                "check": "A_text_vs_image",
                "slide": i + 1,
                "text_title": text_title,
                "matching_image_title": best_match,
                "jaccard": round(best_score, 2),
                "action": "removed"
            })

    return remove, flags


def check_b_image_title_global(prs, pdf_index):
    """
    Keep only the first occurrence of each image title across the full deck.
    Uses the PDF source index for title lookup.
    """
    slides = list(prs.slides)
    seen_titles = {}   # normalised title → first slide num
    remove = set()
    flags = []

    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h not in pdf_index:
            continue
        raw = pdf_index[h][2]
        if not raw:
            continue
        n = normalise(raw)
        if not n:
            continue

        if n in seen_titles:
            remove.add(i)
            flags.append({
                "check": "B_image_title_dup",
                "slide": i + 1,
                "title": raw,
                "first_seen_at": seen_titles[n],
                "action": "removed"
            })
        else:
            seen_titles[n] = i + 1

    return remove, flags


def check_c_exact_md5(prs):
    """Remove exact pixel duplicates."""
    slides = list(prs.slides)
    seen = {}
    remove = set()
    flags = []
    for i, slide in enumerate(slides):
        if not is_fullslide_image(slide, prs):
            continue
        blob = get_img_blob(slide)
        if blob is None:
            continue
        h = hashlib.md5(blob).hexdigest()
        if h in seen:
            remove.add(i)
            flags.append({
                "check": "C_exact_md5",
                "slide": i + 1,
                "first_seen_at": seen[h],
                "action": "removed"
            })
        else:
            seen[h] = i + 1
    return remove, flags


# ── main loop ─────────────────────────────────────────────────────────────────

def quality_check(input_path, output_path=None, dry_run=False):
    if output_path is None:
        output_path = input_path

    print("Building PDF source index (all 5 PDFs)...", file=sys.stderr)
    pdf_index = build_pdf_index()
    print(f"Indexed {len(pdf_index)} PDF pages\n", file=sys.stderr)

    prs = Presentation(input_path)
    total_in = len(list(prs.slides))
    all_flags = []
    iteration = 0

    while True:
        iteration += 1
        iter_flags = []

        # Run all three checks against the CURRENT (unmodified) prs.
        # Collect everything into ONE combined remove set, then apply
        # a SINGLE remove_slides call — this prevents index drift from
        # sequential multi-call removals.
        rc, fc = check_c_exact_md5(prs)
        rb, fb = check_b_image_title_global(prs, pdf_index)
        ra, fa = check_a_text_vs_image(prs, pdf_index)

        iter_flags.extend(fc)
        iter_flags.extend(fb)
        iter_flags.extend(fa)

        # Union of all three remove sets (indices into the CURRENT prs)
        combined = rc | rb | ra
        iter_removed = len(combined)

        print(f"Iteration {iteration}: found {iter_removed} issues "
              f"(C={len(rc)} B={len(rb)} A={len(ra)})", file=sys.stderr)

        if not dry_run and combined:
            remove_slides(prs, combined)

        all_flags.extend(iter_flags)

        # In dry-run mode run exactly one pass — deck is unchanged so no convergence possible
        if dry_run or iter_removed == 0:
            break

    if not dry_run and (all_flags or output_path != input_path):
        prs.save(output_path)

    total_out = len(list(prs.slides))
    result = {
        "input":           input_path,
        "output":          output_path if not dry_run else "(dry-run)",
        "total_in":        total_in,
        "total_out":       total_out,
        "iterations":      iteration,
        "total_removed":   len(all_flags),
        "flags":           all_flags,
        "clean":           len(all_flags) == 0,
    }
    return result


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input",     help="Input PPTX")
    parser.add_argument("output",    nargs="?", help="Output PPTX (omit to overwrite input)")
    parser.add_argument("--dry-run", action="store_true", help="Report only, do not save")
    args = parser.parse_args()

    out = args.output or args.input
    result = quality_check(args.input, out, dry_run=args.dry_run)

    print(json.dumps(result, indent=2))
    status = "CLEAN" if result["clean"] else f"REMOVED {result['total_removed']}"
    mode = " [DRY RUN]" if args.dry_run else ""
    print(f"\nQUALITY CHECK {status}{mode} — {result['total_out']} slides remaining",
          file=sys.stderr)
    sys.exit(0)
