#!/usr/bin/env python3
"""
build_slides.py — all-white (n8n house style) WSQ course deck for TGS-2021010366
'Application Integration with Docker and Kubernetes'.

Design follows the tertiary-course-slides skill (white theme, blue/teal accents, Arial,
logo cover, admin front/back matter, per-lab overview + commands + test, break dividers).
Lab content is driven by labs_data.py so the deck stays 100% aligned with the labs,
Lesson Plan, Learner Guide and Markdown.  Run:  python3 build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import labs_data as L

REPO = os.environ.get("COURSE_REPO") or os.getcwd()
ASSETS = os.path.join(REPO, "courseware", "assets")
# icons ship with the tertiary-course-slides skill
ICONS = None
for cand in [os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "icons"),
             os.path.join(REPO, ".claude/skills/tertiary-course-slides/assets/icons")]:
    if os.path.isdir(cand):
        ICONS = cand; break

TITLE = "Application Integration with Docker and Kubernetes"
CODE = "TGS-2021010366"
VERSION = os.environ.get("DECK_VERSION", "21")
SLUG = "Application-Integration-with-Docker-and-Kubernetes"
FOOT = f"{TITLE}  ·  {CODE}"

NAVY=RGBColor(0x0B,0x12,0x20); BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81)
AMBER=RGBColor(0xF5,0x9E,0x0B); INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72)
LIGHT=RGBColor(0xF5,0xF8,0xFC); WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0)
VIOLET=RGBColor(0x7C,0x3A,0xED); CODEBG=RGBColor(0xF4,0xF7,0xFB)
CENTER=PP_ALIGN.CENTER; MIDDLE=MSO_ANCHOR.MIDDLE; FONT="Arial"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=color
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(1)
    sh.shadow.inherit=False; return sh
def rrect(s,x,y,w,h,fill,line=LINE):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    try: sh.adjustments[0]=0.05
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(1)
    sh.shadow.inherit=False; return sh
def oval(s,x,y,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background(); sh.shadow.inherit=False; return sh
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4,mono=False):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    if runs and not isinstance(runs[0],list): runs=[runs]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0)
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col
            r.font.name=("Consolas" if mono else FONT)
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.margin_left=tf.margin_right=Pt(0)
    first=True
    for it in items:
        lvl=0; text=it
        if isinstance(it,tuple): text,lvl=it
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(gap); p.space_before=Pt(0)
        rm=p.add_run(); rm.text=("   "*lvl)+("•  " if lvl==0 else "–  ")
        rm.font.size=Pt(size); rm.font.color.rgb=(mcolor if lvl==0 else GREY); rm.font.bold=True; rm.font.name=FONT
        r=p.add_run(); r.text=text; r.font.size=Pt(size-lvl); r.font.color.rgb=color; r.font.name=FONT
    return tb
PAGE={"n":0}
def footer(s,dark=False):
    PAGE["n"]+=1; c=WHITE if dark else GREY
    txt(s,Inches(0.55),Inches(7.08),Inches(6.6),Inches(0.3),[[(FOOT,8.5,c,False)]])
    txt(s,Inches(6.0),Inches(7.08),Inches(6),Inches(0.3),[[("© 2026 Tertiary Infotech Academy Pte Ltd",8.5,c,False)]],align=CENTER)
    txt(s,Inches(11.9),Inches(7.08),Inches(0.9),Inches(0.3),[[(str(PAGE["n"]),9,c,True)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,Inches(0.55),Inches(0.60),Inches(0.14),Inches(0.66),kcolor)
    y=Inches(0.55)
    if kicker:
        txt(s,Inches(0.8),Inches(0.5),Inches(11.8),Inches(0.34),[[(kicker,14,kcolor,True)]]); y=Inches(0.82)
    txt(s,Inches(0.8),y,Inches(12.0),Inches(0.85),[[(title,29,INK,True)]])
    rect(s,Inches(0.8),Inches(1.55),Inches(12.0),Pt(2),LINE)
def idisc(s,x,y,d,icon,color=BLUE):
    oval(s,x,y,d,d,color); pad=int(d*0.27)
    p=os.path.join(ICONS or "", f"{icon}.png")
    if ICONS and os.path.exists(p): s.shapes.add_picture(p,x+pad,y+pad,width=d-2*pad,height=d-2*pad)

def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    if os.path.exists(f"{ASSETS}/tertiary-infotech-logo.png"):
        s.shapes.add_picture(f"{ASSETS}/tertiary-infotech-logo.png",Inches(0.85),Inches(0.7),height=Inches(1.05))
    if os.path.exists(f"{ASSETS}/docker-k8s-course-logo.png"):
        s.shapes.add_picture(f"{ASSETS}/docker-k8s-course-logo.png",Inches(11.5),Inches(0.62),height=Inches(1.15))
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[("Application Integration with Docker and Kubernetes",44,INK,True)]])
    rect(s,Inches(0.92),Inches(4.95),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.25),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(11.0),Inches(6.55),Inches(1.85),Inches(0.34),[[(f"Version {VERSION}",13,BLUE,True)]],align=PP_ALIGN.RIGHT)
    txt(s,Inches(0.9),Inches(6.7),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.5),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.3),Inches(0.7),Inches(2.5),Inches(1.6),[[(n,72,LINE,True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=19):
    s=slide(); head(s,title,kicker); bullets(s,Inches(0.85),Inches(1.9),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=slide(); head(s,title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=slide(); head(s,title,kicker)
    xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rrect(s,x,Inches(1.95),Inches(3.65),Inches(4.6),WHITE); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.25),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.8),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
def brk(kind,dur,color):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=CENTER); PAGE["n"]+=1

# ----- admin helpers (white theme) -----
def about_trainer(name=None,role=None,rows=None,blank=False,links=None):
    s=slide(); head(s,"About the Trainer","YOUR TRAINER")
    cx,cy,cw,ch=Inches(0.8),Inches(1.95),Inches(3.7),Inches(4.0)
    rrect(s,cx,cy,cw,ch,INK,line=None)
    d=Inches(1.25); idisc(s,cx+(cw-d)//2,cy+Inches(0.45),d,"ic_grad",BLUE)
    txt(s,cx,cy+Inches(2.0),cw,Inches(0.55),[[(name or "[ Trainer Name ]",20,WHITE,True)]],align=CENTER)
    txt(s,cx,cy+Inches(2.55),cw,Inches(0.45),[[(role or "[ Title / Role ]",14,BLUE,True)]],align=CENTER)
    gap=Inches(0.82) if links else Inches(1.0); rx,ry,dd=Inches(4.85),Inches(2.0),Inches(0.7)
    for icon,label,desc in rows:
        idisc(s,rx,ry,dd,icon,BLUE)
        if blank:
            txt(s,rx+dd+Inches(0.3),ry,Inches(2.3),dd,[[(label,15,INK,True)]],anchor=MIDDLE)
            rect(s,Inches(8.05),ry+dd//2-Pt(1),Inches(4.3),Pt(1.5),INK)
        else:
            txt(s,rx+dd+Inches(0.3),ry-Inches(0.05),Inches(7.5),dd+Inches(0.1),
                [[(label,15,INK,True)],[(desc,13,GREY,False)]],anchor=MIDDLE,space=2)
        ry+=gap
    footer(s)
def digital_attendance():
    s=slide(); head(s,"Digital Attendance (Mandatory)","TRAQOM · SSG DIGITAL ATTENDANCE")
    bullets(s,Inches(0.85),Inches(2.0),Inches(6.3),Inches(4.2),[
        "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
        "The trainer displays the digital attendance QR code from the SSG portal.",
        "Scan the QR code with your phone camera and submit your attendance.",
        "A minimum of 75% attendance is required for assessment and funding."],size=16,gap=14)
    cx,cy,cw,ch=Inches(7.55),Inches(2.0),Inches(4.8),Inches(3.0); rrect(s,cx,cy,cw,ch,LIGHT)
    d=Inches(1.1); idisc(s,cx+(cw-d)//2,cy+Inches(0.35),d,"ic_attend",BLUE)
    txt(s,cx+Inches(0.3),cy+Inches(1.65),cw-Inches(0.6),Inches(1.2),
        [[("Minimum 75%",26,BLUE,True)],[("attendance required",26,BLUE,True)]],align=CENTER,space=2)
    footer(s)
def icon_cards(title,kicker,cards,sub=None,cols=3):
    s=slide(); head(s,title,kicker)
    y0=Inches(2.1) if sub else Inches(1.9)
    if sub: txt(s,Inches(0.85),Inches(1.62),Inches(11.6),Inches(0.34),[[(sub,15,GREY,False)]])
    gap=Inches(0.3); cw=(Inches(11.65)-gap*(cols-1))//cols; ch=Inches(2.45)
    for k,(icon,t,desc) in enumerate(cards):
        col=k%cols; row=k//cols; x=Inches(0.85)+col*(cw+gap); y=y0+row*(ch+Inches(0.28))
        rrect(s,x,y,cw,ch,WHITE); d=Inches(0.92); idisc(s,x+(cw-d)//2,y+Inches(0.32),d,icon,BLUE)
        txt(s,x+Inches(0.2),y+Inches(1.45),cw-Inches(0.4),Inches(0.5),[[(t,17,INK,True)]],align=CENTER)
        txt(s,x+Inches(0.28),y+Inches(1.9),cw-Inches(0.56),Inches(0.5),[[(desc,13,GREY,False)]],align=CENTER)
    footer(s)
def grid_cards(title,kicker,cards,cols=3):
    s=slide(); head(s,title,kicker)
    gap=Inches(0.3); cw=(Inches(11.65)-gap*(cols-1))//cols; ch=Inches(1.95)
    for k,(icon,t,desc) in enumerate(cards):
        col=k%cols; row=k//cols; x=Inches(0.85)+col*(cw+gap); y=Inches(1.9)+row*(ch+Inches(0.27))
        rrect(s,x,y,cw,ch,WHITE); d=Inches(0.62); idisc(s,x+Inches(0.27),y+Inches(0.27),d,icon,BLUE)
        txt(s,x+Inches(0.27)+d+Inches(0.18),y+Inches(0.27),cw-Inches(0.9),d,[[(t,15,INK,True)]],anchor=MIDDLE)
        txt(s,x+Inches(0.27),y+Inches(1.05),cw-Inches(0.54),Inches(0.8),[[(desc,12,GREY,False)]])
    footer(s)
def lesson_plan_cards(days,timing):
    s=slide(); head(s,"Lesson Plan — 2 Days, 8 hours/day","SCHEDULE")
    gap=Inches(0.4); cw=(Inches(11.65)-gap)//2; cy,ch=Inches(1.92),Inches(3.7)
    for k,(num,lines) in enumerate(days):
        x=Inches(0.85)+k*(cw+gap); rrect(s,x,cy,cw,ch,WHITE)
        d=Inches(0.74); ov=oval(s,x+Inches(0.32),cy+Inches(0.32),d,d,BLUE)
        tf=ov.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=CENTER
        r=p.add_run(); r.text=num; r.font.size=Pt(24); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
        txt(s,x+Inches(0.32)+d+Inches(0.24),cy+Inches(0.32),cw-Inches(1.4),d,[[("Day "+num,17,INK,True)]],anchor=MIDDLE)
        para=[[(t,13,INK if bold else GREY,bold)] for t,bold in lines]
        txt(s,x+Inches(0.32),cy+Inches(1.3),cw-Inches(0.64),Inches(2.2),para,space=7)
    rrect(s,Inches(0.85),Inches(5.85),Inches(11.65),Inches(0.62),LIGHT)
    txt(s,Inches(1.2),Inches(5.85),Inches(11),Inches(0.62),[[("Daily timing:   ",14,BLUE,True),(timing,14,INK,False)]],anchor=MIDDLE)
    footer(s)
def learning_outcomes(los):
    s=slide(); head(s,"Learning Outcomes","WHAT YOU'LL ACHIEVE")
    ry,gap,d=Inches(1.85),Inches(0.78),Inches(0.56)
    for i,lo in enumerate(los,1):
        ov=oval(s,Inches(0.85),ry,d,d,BLUE); tf=ov.text_frame; p=tf.paragraphs[0]; p.alignment=CENTER
        r=p.add_run(); r.text=str(i); r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
        txt(s,Inches(0.85)+d+Inches(0.32),ry,Inches(10.8),d,[[("LO"+str(i)+"   ",15,BLUE,True),(lo,15,INK,False)]],anchor=MIDDLE)
        ry+=gap
    footer(s)
def assessment_twocard():
    s=slide(); head(s,"Assessment","FINAL ASSESSMENT")
    cy,ch=Inches(1.95),Inches(4.0); lw=Inches(5.75); lx=Inches(0.85); rx=lx+lw+Inches(0.3)
    rrect(s,lx,cy,lw,ch,WHITE); d=Inches(0.72); idisc(s,lx+Inches(0.34),cy+Inches(0.34),d,"ic_clip",BLUE)
    txt(s,lx+Inches(0.34)+d+Inches(0.24),cy+Inches(0.34),lw-Inches(1.4),d,[[("Final Assessment",18,INK,True)]],anchor=MIDDLE)
    bullets(s,lx+Inches(0.34),cy+Inches(1.4),lw-Inches(0.68),Inches(2.4),[
        "Written / MCQ assessment on the LMS","Held on Day 2 from 4:00 PM",
        "Covers Docker & Kubernetes concepts and the hands-on labs",
        "Open book — slides & Learner Guide","Must be attempted for SSG funding"],size=13,gap=10)
    rrect(s,rx,cy,lw,ch,BLUE,line=None); oval(s,rx+Inches(0.34),cy+Inches(0.34),d,d,RGBColor(0x0B,0x4F,0xC0))
    txt(s,rx+Inches(0.34)+d+Inches(0.24),cy+Inches(0.34),lw-Inches(1.4),d,[[("Funding & Competency",18,WHITE,True)]],anchor=MIDDLE)
    LT=RGBColor(0xCF,0xE0,0xFB)
    txt(s,rx+Inches(0.34),cy+Inches(1.45),lw-Inches(0.68),Inches(1.3),
        [[("Minimum 75% attendance",15,WHITE,True)],[("Based on SSG digital attendance records.",13,LT,False)]],space=6)
    txt(s,rx+Inches(0.34),cy+Inches(2.75),lw-Inches(0.68),Inches(1.3),
        [[("Assessed as Competent",15,WHITE,True)],[("Pass the written / MCQ assessment.",13,LT,False)]],space=6)
    footer(s)

def assessment_flow():
    s=slide(); head(s,"Assessment Flow","ASSESSMENT")
    txt(s,Inches(0.85),Inches(1.6),Inches(11.6),Inches(0.34),[[("At the assessment stage, follow these five steps in order.",15,GREY,False)]])
    steps=[("ic_qr","STEP 1","TRAQOM","Scan the TRAQOM QR code on the LMS and complete the survey."),
           ("ic_attend","STEP 2","Attendance","Take the assessment digital attendance."),
           ("ic_clip","STEP 3","Assessment","Complete the written / MCQ assessment on the LMS."),
           ("ic_upload","STEP 4","Submit Answers","Submit your assessment answers on the LMS."),
           ("ic_sign","STEP 5","Sign the Record","Sign your Assessment Summary Record.")]
    gap=Inches(0.4); cw=(Inches(11.65)-gap*4)//5; cy,ch=Inches(2.35),Inches(2.7)
    for k,(icon,step,t,desc) in enumerate(steps):
        x=Inches(0.85)+k*(cw+gap); rrect(s,x,cy,cw,ch,WHITE if k%2==0 else LIGHT)
        d=Inches(0.76); idisc(s,x+(cw-d)//2,cy+Inches(0.25),d,icon,BLUE)
        txt(s,x+Inches(0.1),cy+Inches(1.15),cw-Inches(0.2),Inches(0.3),[[(step,10,BLUE,True)]],align=CENTER)
        txt(s,x+Inches(0.08),cy+Inches(1.42),cw-Inches(0.16),Inches(0.6),[[(t,13,INK,True)]],align=CENTER)
        txt(s,x+Inches(0.12),cy+Inches(2.0),cw-Inches(0.24),Inches(0.6),[[(desc,9.5,GREY,False)]],align=CENTER)
        if k<4: txt(s,x+cw+Inches(0.04),cy+Inches(1.0),Inches(0.32),Inches(0.5),[[("›",28,BLUE,True)]],align=CENTER)
    txt(s,Inches(0.85),Inches(5.5),Inches(11.6),Inches(0.34),
        [[("TIP   ",13,TEAL,True),("Courseware & the assessment are on the LMS:  https://lms.tertiaryinfotech.com",13,GREY,False)]])
    footer(s)

def cert_traqom():
    s=slide(); head(s,"Certification & TRAQOM Survey","BEFORE YOU GO")
    cy,ch=Inches(2.0),Inches(3.6); lw=Inches(5.75); lx=Inches(0.85); rx=lx+lw+Inches(0.3)
    rrect(s,lx,cy,lw,ch,WHITE)
    d=Inches(0.78); idisc(s,lx+Inches(0.36),cy+Inches(0.36),d,"ic_survey",BLUE)
    txt(s,lx+Inches(0.36)+d+Inches(0.26),cy+Inches(0.36),lw-Inches(1.5),d,[[("TRAQOM Survey (Mandatory)",17,INK,True)]],anchor=MIDDLE)
    txt(s,lx+Inches(0.36),cy+Inches(1.5),lw-Inches(0.72),Inches(1.6),
        [[("Complete the certification & TRAQOM survey to receive your certificate.",14.5,INK,False)],
         [("lms.tertiaryinfotech.com",14.5,BLUE,True)]],space=8)
    rrect(s,rx,cy,lw,ch,BLUE,line=None)
    oval(s,rx+Inches(0.36),cy+Inches(0.36),d,d,RGBColor(0x0B,0x4F,0xC0))
    icp=os.path.join(ICONS or "","ic_check.png")
    if ICONS and os.path.exists(icp):
        pad=int(d*0.27); s.shapes.add_picture(icp,rx+Inches(0.36)+pad,cy+Inches(0.36)+pad,width=d-2*pad,height=d-2*pad)
    txt(s,rx+Inches(0.36)+d+Inches(0.26),cy+Inches(0.36),lw-Inches(1.5),d,[[("Digital Attendance",17,WHITE,True)]],anchor=MIDDLE)
    txt(s,rx+Inches(0.36),cy+Inches(1.5),lw-Inches(0.72),Inches(1.6),
        [[("Remember to take AM, PM and Assessment digital attendance — required for your funding.",14.5,RGBColor(0xCF,0xE0,0xFB),False)]])
    footer(s)

# ----- code + lab slides (driven by labs_data) -----
def code_card(s,x,y,w,cd,accent=TEAL):
    n=cd.count("\n")+1; ch=Inches(0.16+0.205*n)
    rrect(s,x,y,w,ch,CODEBG); rect(s,x,y,Inches(0.09),ch,accent)
    tb=s.shapes.add_textbox(x+Inches(0.24),y+Inches(0.07),w-Inches(0.36),ch-Inches(0.1))
    tf=tb.text_frame; tf.word_wrap=True; tf.margin_left=tf.margin_right=Pt(0); tf.margin_top=tf.margin_bottom=Pt(0)
    for i,ln in enumerate(cd.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(0); p.space_before=Pt(0)
        r=p.add_run(); r.text=ln or " "; r.font.size=Pt(10.5); r.font.name="Consolas"; r.font.color.rgb=INK
    return ch

import re as _re
def clean(t):
    """Strip light markdown (** and `) so it doesn't show literally on slides."""
    return t.replace("**","").replace("`","")

def _seg_height(cap,cd):
    if cap=="Note":
        n=max(1,len(cd)//95+1)
        return Inches(0.34*n+0.3)
    n=cd.count("\n")+1
    return (Inches(0.42) if cap else Inches(0)) + Inches(0.16+0.205*n) + Inches(0.2)

def lab_segments(lab):
    """Convert a lab's body into (caption, code) pairs for command slides."""
    segs=[]; cap=None
    for blk in lab["body"]:
        k=blk[0]
        if k in ("h3","p"):
            t=blk[1]
            if k=="h3": cap=t
            elif cap is None: cap=t
        elif k=="code":
            segs.append((cap, blk[1])); cap=None
        elif k=="note":
            segs.append(("Note", "» "+blk[1])); cap=None
        # tables skipped on slides (present in the Learner Guide)
    return segs

def cmd_name(cd):
    for ln in cd.split("\n"):
        ln=ln.strip()
        if ln and not ln.startswith("#"):
            return " ".join(ln.split()[:2])
    return ""

def lab_overview(lab,kicker):
    s=slide(); head(s,lab["title"],kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.82),Inches(1.85),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.87),Inches(1.85),Inches(0.4),[[(f"LAB {lab['num']}",16,WHITE,True)]],align=CENTER)
    txt(s,Inches(2.95),Inches(1.85),Inches(9.6),Inches(0.45),[[(lab["topic"],14,GREY,True)]],anchor=MIDDLE)
    txt(s,Inches(0.85),Inches(2.5),Inches(11.7),Inches(1.7),[[(clean(lab["goal"]),19,INK,False)]])
    rrect(s,Inches(0.85),Inches(4.45),Inches(11.7),Inches(1.95),LIGHT)
    txt(s,Inches(1.1),Inches(4.65),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(5.05),Inches(11.2),Inches(0.7),[[(clean(lab["build"]),17,INK,True)]])
    cmds=" · ".join(dict.fromkeys(cmd_name(c) for _,c in lab_segments(lab) if cmd_name(c)))[:130]
    txt(s,Inches(1.1),Inches(5.85),Inches(11),Inches(0.5),[[("Key commands:  ",13,GREY,True),(cmds,13,GREY,False)]])
    footer(s)

def lab_command_slides(lab,kicker):
    segs=lab_segments(lab); budget=Inches(4.95); page=[]; used=Inches(0)
    chunks=[]
    for seg in segs:
        h=_seg_height(*seg)
        if page and used+h>budget:
            chunks.append(page); page=[]; used=Inches(0)
        page.append(seg); used+=h
    if page: chunks.append(page)
    for ci,chunk in enumerate(chunks):
        s=slide(); ttl=lab["title"]+(" (cont.)" if ci else "")
        head(s,ttl,f"LAB {lab['num']}  ·  COMMANDS",TEAL)
        y=Inches(1.85)
        for cap,cd in chunk:
            if cap=="Note":
                body=clean(cd[2:] if cd.startswith("» ") else cd)
                txt(s,Inches(0.85),y,Inches(11.65),Inches(0.34),[[("Note   ",13,AMBER,True),(body,13,GREY,False)]])
                y+=_seg_height("Note",cd); continue
            if cap:
                txt(s,Inches(0.85),y,Inches(11.6),Inches(0.34),[[(clean(cap),14,BLUE,True)]]); y+=Inches(0.42)
            ch=code_card(s,Inches(0.85),y,Inches(11.65),cd,TEAL); y+=ch+Inches(0.2)
        footer(s)

def test_slide(lab):
    s=slide(); head(s,lab["title"],f"LAB {lab['num']}  ·  TEST IT",TEAL)
    rrect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.8),RGBColor(0xE8,0xF7,0xEE),line=None)
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",22,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11.1),Inches(1.6),[[(clean(lab["test"]),18,INK,False)]])
    txt(s,Inches(0.85),Inches(5.45),Inches(11.7),Inches(0.4),
        [[("KillerCoda:  ",13,TEAL,True),(lab["killercoda"],13,GREY,False)]])
    footer(s)

def lab_deck(num,kicker):
    lab=L.lab_by_num(num)
    lab_overview(lab,kicker); lab_command_slides(lab,kicker); test_slide(lab)

from build_slides_content import build   # concept slides + sequence
build(globals())

OUT=os.path.join(REPO,"courseware",f"{SLUG}-v{VERSION}.pptx")
prs.save(OUT)
print("Saved deck:", os.path.basename(OUT), "|", PAGE["n"], "numbered slides; total", len(prs.slides.__iter__.__self__._sldIdLst))
